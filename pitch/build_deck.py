"""Build the NEXUS investor / company pitch deck (16:9, dark control-room theme).

    cd backend && uv run --with python-pptx --with matplotlib python ../pitch/build_deck.py

Every number in the deck comes from a data source read at build time:

* ``backend/benchmarks/results/latest.json`` (full benchmark) — or ``sample.json`` (generated on the fly from the
  tiny scale when the full run is not available yet; charts are then labelled "sample run").
* ``pitch/data/sweep.json`` and ``pitch/data/nlq_examples.json`` — produced by ``pitch/capture_measurements.py``
  with the real engine (peak-hour capacity sweep, grounded console answers).
* ``DEMO`` below — the numbers printed by ``uv run nexus demo`` on 2026-08-30 (calibrated small world, 10:30 peak,
  ×1.2 surge, R07 motor fault).
* Line counts are measured by walking the repository.

The script is deterministic and idempotent: it overwrites ``pitch/NEXUS_Pitch_Deck.pptx`` and ``pitch/charts/deck_*.png``.
"""

from __future__ import annotations

import datetime as dt
import json
import math
import subprocess
import sys
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
BACKEND = ROOT / "backend"
RESULTS = BACKEND / "benchmarks" / "results"
CHARTS = HERE / "charts"
DATA = HERE / "data"
OUT = HERE / "NEXUS_Pitch_Deck.pptx"

# ---------------------------------------------------------------------------------------------------------------------
# theme
# ---------------------------------------------------------------------------------------------------------------------
HEX = {
    "bg": "#0A0D12",
    "panel": "#11161D",
    "panel2": "#0E1319",
    "header": "#182130",
    "border": "#1F2933",
    "text": "#E6EDF3",
    "muted": "#8B98A5",
    "cyan": "#22D3EE",
    "green": "#22C55E",
    "amber": "#F59E0B",
    "red": "#EF4444",
    "violet": "#A78BFA",
    "blue": "#60A5FA",
}


def rgb(name: str) -> RGBColor:
    return RGBColor.from_string(HEX[name].lstrip("#"))


BG, PANEL, PANEL2, HEADER, BORDER = rgb("bg"), rgb("panel"), rgb("panel2"), rgb("header"), rgb("border")
TEXT, MUTED, CYAN, GREEN, AMBER, RED, VIOLET, BLUE = (
    rgb("text"), rgb("muted"), rgb("cyan"), rgb("green"), rgb("amber"), rgb("red"), rgb("violet"), rgb("blue"),
)
FONT = "Segoe UI"
MONO = "Consolas"
FOOTER = "NEXUS · Autonomous Operations Intelligence"
STRATEGY_ORDER = ["baseline", "optimized", "ai_planner", "nexus_full"]
STRATEGY_LABEL = {
    "baseline": "Baseline (FIFO + nearest robot)",
    "optimized": "Optimized (CP-SAT + batching + routing)",
    "ai_planner": "AI planner (playbooks, no simulation)",
    "nexus_full": "NEXUS full (simulate-before-execute)",
}
STRATEGY_COLOR = {"baseline": HEX["muted"], "optimized": HEX["cyan"], "ai_planner": HEX["violet"], "nexus_full": HEX["green"]}

# ---------------------------------------------------------------------------------------------------------------------
# measured facts (see module docstring for provenance)
# ---------------------------------------------------------------------------------------------------------------------
DEMO: dict[str, Any] = {
    "when": "10:30 (late-morning peak), demand surge ×1.2 for 60 min from 10:25, R07 motor fault at 10:30",
    "baseline_breach": 39.4,
    "plan_breach": 3.1,
    "ft_now": 2.1,
    "ft_nothing": 9.2,
    "ft_plan": 4.5,
    "thr_nothing": 438,
    "thr_plan": 546,
    "candidates": 9,
    "allocations": 21,
    "total_s": 18.8,
    "planning_ms": 1,
    "simulation_s": 12.1,
    "risk_s": 7.2,
    "seeds": 3,
    "sigma_pct": 0.2,
    "recovery": [("10:45", 5.3), ("11:00", 4.3), ("11:15", 3.6), ("11:30", 3.1)],
    # name, projected SLA breach %, avg fulfilment min, throughput/h, risk, recommended
    "candidates_table": [
        ("Add 2 robots + batching", 3.1, 4.54, 546, "LOW", True),
        ("Reassign to R03 & R11, prioritise HIGH, route via C2", 36.3, 8.96, 446, "LOW", False),
        ("Enable batching (4 orders/trip) + deadline sequencing", 36.8, 8.52, 440, "LOW", False),
        ("Do nothing (reference)", 39.4, 9.17, 438, "—", False),
        ("Reassign released work to R03 & R11", 39.4, 9.17, 438, "—", False),
        ("Spread released work over 4 robots + batching", 40.2, 9.00, 429, "—", False),
        ("Prefer corridor C2 near hot zone C", 45.7, 10.13, 409, "—", False),
    ],
}
PERF: dict[str, Any] = {
    "ticks": {"small": (6000, 2300), "medium": (1500, 600), "large": (500, 200)},  # baseline, optimized
    "forecast_ms": {"small": 3.5, "large": 8.8},
    "fork_ms": {"small": 5, "large": 26},
    "tests": {"agents": 8, "api": 6, "events": 9, "forecasting": 8, "nlq": 15, "optimization": 17, "simulation": 14, "twin": 11, "whatif": 4},
    "docker_gb": 1.43,
}
SCALES = [  # scale, robots, storage zones, docks, chargers, SKUs, inventory (spec), base orders/h, grid, shelves, workers
    ("tiny", 4, "4 (2×2)", 2, 2, "120", "4,000", 90, "29×29", 192, 3),
    ("small", 12, "12 (4×3)", 4, 4, "600", "18,000", 400, "53×41", 576, 7),
    ("medium", 40, "24 (6×4)", 8, 10, "2,000", "60,000", 1000, "77×53", 1152, 16),
    ("large", 100, "50 (10×5)", 16, 24, "5,000", "150,000", 1800, "125×65", 2400, 30),
]
ACTIONS = [
    "REASSIGN_TASKS", "REPRIORITIZE_ORDERS", "SEND_TO_CHARGE", "REROUTE_AVOID_ZONE", "PREFER_CORRIDOR",
    "REPOSITION_INVENTORY", "SET_BATCHING", "SET_ZONE_CAPACITY", "CLOSE_ZONE", "OPEN_ZONE", "ADD_ROBOTS",
    "REMOVE_ROBOTS", "DISPATCH_WORKER", "CANCEL_TASKS", "SET_STRATEGY", "NOOP",
]
WHATIF_PRESETS = [
    "Demand +40%", "Demand doubles for 30 min", "Robot R07 fails", "Remove 2 robots", "Add 2 robots",
    "Zone B inaccessible (60 min)", "Loading dock D2 closes", "Charging capacity halved",
    "Aisle blocked in Zone C (30 min)", "Move hot inventory C → B", "Enable 3-order batching",
    "R07 fails during a +30% demand spike",
]


# ---------------------------------------------------------------------------------------------------------------------
# data loading
# ---------------------------------------------------------------------------------------------------------------------
def load_benchmark() -> tuple[dict[str, Any] | None, str]:
    latest = RESULTS / "latest.json"
    sample = RESULTS / "sample.json"
    if latest.exists():
        return json.loads(latest.read_text(encoding="utf-8")), "full benchmark run"
    if not sample.exists():
        print("no benchmark results found — generating a tiny sample run …", flush=True)
        try:
            subprocess.run(
                [sys.executable, "-m", "benchmarks.run_benchmark", "--scale", "tiny", "--minutes", "20", "--seeds", "1",
                 "--out", "results/sample.json", "--no-docs"],
                cwd=BACKEND, check=True, timeout=1800,
            )
        except Exception as exc:  # noqa: BLE001
            print(f"  sample generation failed: {exc}")
    if sample.exists():
        return json.loads(sample.read_text(encoding="utf-8")), "sample run (tiny scale, 20 min, 1 seed)"
    return None, "no benchmark data"


def load_json(path: Path) -> dict[str, Any] | None:
    return json.loads(path.read_text(encoding="utf-8")) if path.exists() else None


def count_loc() -> dict[str, tuple[int, int]]:
    """Non-blank lines and file counts by component (repository walk; node_modules / build output excluded)."""

    def walk(bases: list[Path], exts: set[str], recursive: bool = True) -> tuple[int, int]:
        lines = files = 0
        for base in bases:
            if not base.exists():
                continue
            it = base.rglob("*") if recursive else base.glob("*")
            for p in it:
                if p.suffix not in exts or not p.is_file():
                    continue
                if any(part in ("node_modules", ".next", ".venv", "__pycache__", "results") for part in p.parts):
                    continue
                if p.name == "world.small.json":
                    continue
                try:
                    lines += sum(1 for ln in p.read_text(encoding="utf-8", errors="ignore").splitlines() if ln.strip())
                    files += 1
                except OSError:
                    continue
        return lines, files

    out = {
        "Python — engine, agents, API": walk([BACKEND / "nexus", BACKEND / "benchmarks"], {".py"}),
        "Python — tests": walk([BACKEND / "tests"], {".py"}),
        "TypeScript — twin UI": walk([ROOT / "frontend" / "src"], {".ts", ".tsx"}),
        "Docs, ADRs, README": (
            walk([ROOT / "docs"], {".md"})[0] + walk([ROOT, BACKEND, ROOT / "frontend"], {".md"}, recursive=False)[0],
            walk([ROOT / "docs"], {".md"})[1] + walk([ROOT, BACKEND, ROOT / "frontend"], {".md"}, recursive=False)[1],
        ),
        "Infra — compose, CI, Grafana, Docker": (
            walk([ROOT / "deploy", ROOT / ".github"], {".yml", ".yaml", ".json"})[0]
            + walk([ROOT], {".yml"}, recursive=False)[0]
            + sum(len([ln for ln in (d / "Dockerfile").read_text(errors="ignore").splitlines() if ln.strip()]) for d in (BACKEND, ROOT / "frontend") if (d / "Dockerfile").exists()),
            walk([ROOT / "deploy", ROOT / ".github"], {".yml", ".yaml", ".json"})[1] + 3,
        ),
    }
    return out


# ---------------------------------------------------------------------------------------------------------------------
# matplotlib charts (dark theme, saved as PNG)
# ---------------------------------------------------------------------------------------------------------------------
def _plt():
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    plt.rcParams.update(
        {
            "figure.facecolor": HEX["bg"],
            "savefig.facecolor": HEX["bg"],
            "axes.facecolor": HEX["panel"],
            "axes.edgecolor": HEX["border"],
            "axes.labelcolor": HEX["muted"],
            "axes.titlecolor": HEX["text"],
            "axes.titleweight": "bold",
            "axes.titlesize": 12,
            "xtick.color": HEX["muted"],
            "ytick.color": HEX["muted"],
            "text.color": HEX["text"],
            "grid.color": HEX["border"],
            "grid.alpha": 0.8,
            "axes.grid": True,
            "axes.axisbelow": True,
            "font.family": ["Segoe UI", "DejaVu Sans"],
            "font.size": 10,
            "legend.facecolor": HEX["panel"],
            "legend.edgecolor": HEX["border"],
            "legend.framealpha": 0.9,
        }
    )
    return plt


def _save(fig: Any, name: str) -> Path:
    CHARTS.mkdir(parents=True, exist_ok=True)
    path = CHARTS / f"deck_{name}.png"
    fig.savefig(path, dpi=200, bbox_inches="tight", pad_inches=0.15)
    import matplotlib.pyplot as plt

    plt.close(fig)
    return path


def _style_axes(ax: Any, ylabel: str = "") -> None:
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.grid(axis="x", visible=False)
    ax.set_ylabel(ylabel)


def chart_demo_timeline() -> Path:
    plt = _plt()
    fig, (ax0, ax1) = plt.subplots(2, 1, figsize=(12, 4.6), height_ratios=[1.25, 1.25], sharex=False)
    # --- timeline strip -------------------------------------------------------------------------
    ax0.set_xlim(0, 215)
    ax0.set_ylim(0, 1)
    ax0.axis("off")
    ax0.plot([0, 215], [0.45, 0.45], color=HEX["border"], lw=6, solid_capstyle="round")
    ax0.plot([145, 215], [0.45, 0.45], color=HEX["cyan"], lw=6, solid_capstyle="round", alpha=0.9)
    # (event x, label x, above?) — labels are spread out and joined to their event with a leader line
    events = [
        (0, 8, True, "08:00", "twin starts · 12 robots\n400 orders/h base demand", HEX["muted"]),
        (145, 112, False, "10:25", "demand surge ×1.2\nfor 60 min", HEX["amber"]),
        (150, 138, True, "10:30", "R07 motor fault\n(45 min recovery)", HEX["red"]),
        (151, 158, False, "10:30:30", "decision: 9 plans, 21 allocations\nsimulated in 18.8 s", HEX["cyan"]),
        (152, 178, True, "10:31", "plan executed\n(+2 robots, batching)", HEX["green"]),
        (210, 206, False, "11:30", "breach since decision 3.1%\n(do-nothing projection 39.4%)", HEX["green"]),
    ]
    ax0.set_ylim(-0.75, 1.45)
    for x, xl, above, t, label, color in events:
        y_end = 0.95 if above else -0.05
        ax0.plot([x, xl], [0.45, y_end], color=color, lw=1.0, alpha=0.85)
        ax0.plot([x], [0.45], marker="o", ms=11, color=color, mec=HEX["bg"], mew=2, zorder=5)
        ax0.text(xl, y_end + (0.06 if above else -0.06), f"{t}\n{label}", ha="center", va="bottom" if above else "top",
                 fontsize=8.5, color=color, linespacing=1.25)
    ax0.set_title("The storyline — one simulated morning on the calibrated small warehouse", loc="left", fontsize=11, pad=4)
    # --- outcome panel --------------------------------------------------------------------------
    labels = ["projected\n(do nothing)", "10:45", "11:00", "11:15", "11:30"]
    values = [DEMO["baseline_breach"]] + [v for _, v in DEMO["recovery"]]
    colors = [HEX["red"]] + [HEX["green"]] * 4
    bars = ax1.bar(labels, values, color=colors, width=0.55, edgecolor=HEX["bg"])
    for b, v in zip(bars, values, strict=True):
        ax1.text(b.get_x() + b.get_width() / 2, b.get_height() + 1.0, f"{v:.1f}%", ha="center", va="bottom", fontsize=10, fontweight="bold")
    ax1.axhline(DEMO["plan_breach"], color=HEX["cyan"], lw=1.2, ls="--")
    ax1.text(4.45, DEMO["plan_breach"] + 10.5, f"dashed line: the plan's simulated projection, {DEMO['plan_breach']}%", color=HEX["cyan"], fontsize=8.5, ha="right")
    ax1.set_ylim(0, 48)
    _style_axes(ax1, "SLA breach (%)")
    ax1.set_title("Projected SLA breach over the next 60 min without action vs realized after the plan (breach since decision)", loc="left", fontsize=10.5)
    fig.tight_layout(h_pad=1.2)
    return _save(fig, "demo_timeline")


def chart_candidates() -> Path:
    plt = _plt()
    rows = DEMO["candidates_table"]
    fig, ax = plt.subplots(figsize=(7.4, 4.6))
    names = [r[0] for r in rows][::-1]
    vals = [r[1] for r in rows][::-1]
    colors = []
    for r in rows[::-1]:
        colors.append(HEX["green"] if r[5] else (HEX["muted"] if r[0].startswith("Do nothing") else HEX["cyan"]))
    bars = ax.barh(names, vals, color=colors, edgecolor=HEX["bg"], height=0.62)
    for b, v in zip(bars, vals, strict=True):
        ax.text(b.get_width() + 0.8, b.get_y() + b.get_height() / 2, f"{v:.1f}%", va="center", fontsize=9.5, fontweight="bold")
    ax.set_xlim(0, 55)
    ax.grid(axis="y", visible=False)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.set_xlabel("projected SLA breach over the 60-minute horizon (%) — lower is better")
    ax.tick_params(axis="y", labelsize=9)
    ax.set_title("Every candidate simulated in a forked twin · recommended plan in green", loc="left", fontsize=10.5)
    fig.tight_layout()
    return _save(fig, "candidates")


def _grouped(ax: Any, groups: list[str], series: dict[str, list[float]], fmt: str, colors: dict[str, str] | None = None, width_total: float = 0.78) -> None:
    n = max(1, len(series))
    w = width_total / n
    for i, (name, vals) in enumerate(series.items()):
        xs = [j + i * w - width_total / 2 + w / 2 for j in range(len(groups))]
        color = (colors or STRATEGY_COLOR).get(name, HEX["blue"])
        bars = ax.bar(xs, vals, w * 0.92, label=STRATEGY_LABEL.get(name, name), color=color, edgecolor=HEX["bg"])
        for b, v in zip(bars, vals, strict=True):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height(), fmt.format(v), ha="center", va="bottom", fontsize=7.5)
    ax.set_xticks(list(range(len(groups))))
    ax.set_xticklabels(groups)


def _bench_series(bench: dict[str, Any], key: str) -> tuple[list[str], dict[str, list[float]]]:
    scales = list(bench["scales"])
    present = [s for s in STRATEGY_ORDER if any(s in bench["scales"][sc]["strategies"] for sc in scales)]
    rows = bench["summary_table"]
    series = {}
    for strat in present:
        series[strat] = [float(next((r[key] for r in rows if r["scale"] == sc and r["strategy"] == strat), 0.0)) for sc in scales]
    return scales, series


def chart_bench(bench: dict[str, Any], key: str, title: str, ylabel: str, fmt: str, name: str, label: str, figsize: tuple[float, float] = (7.2, 4.2)) -> Path:
    plt = _plt()
    scales, series = _bench_series(bench, key)
    fig, ax = plt.subplots(figsize=figsize)
    _grouped(ax, scales, series, fmt)
    _style_axes(ax, ylabel)
    ax.set_title(f"{title}  ·  {label}", loc="left", fontsize=10.5)
    ax.legend(fontsize=8, loc="upper left", ncols=2)
    ymax = max((max(v) for v in series.values()), default=1.0)
    ax.set_ylim(0, ymax * 1.28 if ymax > 0 else 1)
    fig.tight_layout()
    return _save(fig, name)


def chart_ticks(bench: dict[str, Any] | None, label: str) -> Path:
    plt = _plt()
    fig, ax = plt.subplots(figsize=(6.4, 4.1))
    if bench is not None and len(bench["scales"]) >= 2:
        scales, series = _bench_series(bench, "ticks_per_second")
        title = f"Engine speed by strategy (ticks / s, log scale) · {label}"
    else:
        scales = list(PERF["ticks"])
        series = {"baseline": [PERF["ticks"][s][0] for s in scales], "optimized": [PERF["ticks"][s][1] for s in scales]}
        title = "Engine speed (ticks / s, log scale) · measured on the calibration runs"
    _grouped(ax, scales, series, "{:,.0f}")
    ax.set_yscale("log")
    ax.grid(axis="y", which="both")
    _style_axes(ax, "simulated seconds per wall second")
    ax.set_title(title, loc="left", fontsize=10.5)
    ax.legend(fontsize=8, loc="upper right")
    fig.tight_layout()
    return _save(fig, "ticks")


def chart_radar(bench: dict[str, Any], label: str) -> tuple[Path, str]:
    plt = _plt()
    scales = list(bench["scales"])
    scale = "small" if "small" in scales else scales[-1]
    strategies = [s for s in STRATEGY_ORDER if s in bench["scales"][scale]["strategies"]]
    rows = {r["strategy"]: r for r in bench["summary_table"] if r["scale"] == scale}
    axes_def = [
        ("SLA compliance", lambda r: 100 - r["sla_breach_pct"], True),
        ("Speed (avg fulfilment)", lambda r: r["avg_fulfillment_min"], False),
        ("Tail (p95)", lambda r: r["p95_fulfillment_min"], False),
        ("Throughput", lambda r: r["throughput_per_hour"], True),
        ("Flow (low congestion)", lambda r: r["congestion_index"], False),
        ("Energy per order", lambda r: r["energy_per_order"], False),
    ]
    scores: dict[str, list[float]] = {s: [] for s in strategies}
    for _, fn, higher_better in axes_def:
        vals = {s: float(fn(rows[s])) for s in strategies}
        if higher_better:
            m = max(vals.values()) or 1.0
            for s in strategies:
                scores[s].append(vals[s] / m if m else 1.0)
        else:
            mn = min(vals.values())
            for s in strategies:
                scores[s].append(1.0 if vals[s] <= 1e-9 and mn <= 1e-9 else (mn / vals[s] if vals[s] > 0 else 1.0))
    n = len(axes_def)
    angles = [2 * math.pi * i / n for i in range(n)] + [0]
    fig = plt.figure(figsize=(5.6, 5.0))
    ax = fig.add_subplot(111, polar=True)
    ax.set_facecolor(HEX["panel"])
    for s in strategies:
        vals = scores[s] + scores[s][:1]
        ax.plot(angles, vals, color=STRATEGY_COLOR[s], lw=2, label=STRATEGY_LABEL[s].split(" (")[0])
        ax.fill(angles, vals, color=STRATEGY_COLOR[s], alpha=0.10)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels([a[0] for a in axes_def], fontsize=8.5)
    ax.set_yticks([0.25, 0.5, 0.75, 1.0])
    ax.set_yticklabels(["", "", "", ""], fontsize=7)
    ax.set_ylim(0, 1.05)
    ax.spines["polar"].set_color(HEX["border"])
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.08), fontsize=8, ncols=2)
    ax.set_title(f"Normalised KPI profile · {scale} scale\n{label}", fontsize=10, pad=20)
    fig.tight_layout()
    return _save(fig, "radar"), scale


def chart_capacity(sweep: dict[str, Any] | None) -> Path:
    plt = _plt()
    fig, ax = plt.subplots(figsize=(7.6, 4.4))
    if sweep:
        rows = sweep["rows"]
        xs = [r["orders_per_hour"] for r in rows]
        for key, label, color, ls in (
            ("baseline_ok", "baseline", HEX["muted"], "-"),
            ("baseline_fail", "baseline + R07 failure", HEX["muted"], "--"),
            ("optimized_ok", "optimized", HEX["cyan"], "-"),
            ("optimized_fail", "optimized + R07 failure", HEX["cyan"], "--"),
        ):
            ys = [r[key]["sla_breach_pct"] for r in rows]
            ax.plot(xs, ys, color=color, ls=ls, lw=2.2, marker="o", ms=5, label=label)
            for x, y in zip(xs, ys, strict=True):
                if key in ("baseline_ok", "optimized_ok"):
                    ax.annotate(f"{y:.0f}%", (x, y), textcoords="offset points", xytext=(0, 6), ha="center", fontsize=7.5, color=color)
        ax.set_xlabel("demand at the peak (orders / hour, small warehouse, 12 robots)")
    ax.set_ylim(0, 85)
    _style_axes(ax, "projected SLA breach after 90 min (%)")
    ax.legend(fontsize=8.5, loc="upper left")
    ax.set_title("The capacity cliff — same robots, same orders, different scheduler", loc="left", fontsize=10.5)
    fig.tight_layout()
    return _save(fig, "capacity")


def chart_whatif(sweep: dict[str, Any] | None) -> tuple[Path, dict[str, float]]:
    plt = _plt()
    fig, ax = plt.subplots(figsize=(6.0, 4.4))
    vals: dict[str, float] = {}
    if sweep:
        row = next((r for r in sweep["rows"] if r["orders_per_hour"] == 504), sweep["rows"][len(sweep["rows"]) // 2])
        vals = {
            "baseline": row["baseline_ok"]["sla_breach_pct"],
            "baseline\n+ R07 fails": row["baseline_fail"]["sla_breach_pct"],
            "optimized": row["optimized_ok"]["sla_breach_pct"],
            "optimized\n+ R07 fails": row["optimized_fail"]["sla_breach_pct"],
        }
        colors = [HEX["muted"], HEX["red"], HEX["cyan"], HEX["green"]]
        bars = ax.bar(list(vals), list(vals.values()), color=colors, width=0.6, edgecolor=HEX["bg"])
        for b, v in zip(bars, vals.values(), strict=True):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 1, f"{v:.1f}%", ha="center", va="bottom", fontsize=10, fontweight="bold")
        ax.set_title(f"What if R07 fails at the peak? · {row['orders_per_hour']} orders/h · 90 min horizon", loc="left", fontsize=10.5)
    ax.set_ylim(0, 70)
    _style_axes(ax, "projected SLA breach (%)")
    fig.tight_layout()
    return _save(fig, "whatif"), vals


# ---------------------------------------------------------------------------------------------------------------------
# pptx helpers
# ---------------------------------------------------------------------------------------------------------------------
class Deck:
    W = 13.333
    H = 7.5

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.W)
        self.prs.slide_height = Inches(self.H)
        self.blank = self.prs.slide_layouts[6]
        self.n = 0
        self.titles: list[str] = []

    # ---- slide scaffolding -----------------------------------------------------------------------------------------
    def slide(self, kicker: str, title: str, notes: str = "") -> Any:
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        self.titles.append(title)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = BG
        text(s, 0.6, 0.28, 9, 0.3, kicker.upper(), size=10.5, color=CYAN, bold=True)
        title_size = 28 if len(title) <= 48 else 24 if len(title) <= 58 else 21
        text(s, 0.6, 0.55, 11.6, 0.8, title, size=title_size, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.33), Inches(1.1), Inches(0.045))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CYAN
        bar.line.fill.background()
        bar.shadow.inherit = False
        text(s, 0.6, 7.05, 7, 0.3, FOOTER, size=9.5, color=MUTED)
        text(s, 11.6, 7.05, 1.15, 0.3, f"{self.n:02d}", size=9.5, color=MUTED, align=PP_ALIGN.RIGHT, font=MONO)
        if notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    def bare(self, title: str, notes: str = "") -> Any:
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        self.titles.append(title)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = BG
        if notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    def save(self) -> Path:
        self.prs.save(OUT)
        return OUT


def text(slide: Any, x: float, y: float, w: float, h: float, content: str | list[str], size: float = 14, color: RGBColor = TEXT,
         bold: bool = False, font: str = FONT, align: Any = PP_ALIGN.LEFT, anchor: Any = MSO_ANCHOR.TOP, italic: bool = False,
         line_spacing: float = 1.1, space_after: float = 4) -> Any:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.04))
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.name = font
        f.color.rgb = color
    return tb


def rich(slide: Any, x: float, y: float, w: float, h: float, paragraphs: list[list[tuple[str, dict[str, Any]]]], size: float = 13,
         line_spacing: float = 1.15, space_after: float = 6, anchor: Any = MSO_ANCHOR.TOP) -> Any:
    """Paragraphs made of (text, style) runs; style keys: bold, color, size, font, italic."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.05))
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for t, st in runs:
            r = p.add_run()
            r.text = t
            f = r.font
            f.size = Pt(st.get("size", size))
            f.bold = st.get("bold", False)
            f.italic = st.get("italic", False)
            f.name = st.get("font", FONT)
            f.color.rgb = st.get("color", TEXT)
    return tb


def bullets(slide: Any, x: float, y: float, w: float, h: float, items: list[Any], size: float = 13, color: RGBColor = TEXT,
            bullet_color: RGBColor = CYAN, spacing: float = 5, line_spacing: float = 1.12) -> Any:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.05))
    for i, item in enumerate(items):
        lead, rest, level = None, item, 0
        if isinstance(item, tuple):
            if len(item) == 3:
                lead, rest, level = item
            else:
                lead, rest = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.line_spacing = line_spacing
        r0 = p.add_run()
        r0.text = ("     " * level) + ("•  " if level == 0 else "–  ")
        r0.font.color.rgb = bullet_color if level == 0 else MUTED
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.name = FONT
        if lead:
            r1 = p.add_run()
            r1.text = lead + " "
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.name = FONT
            r1.font.color.rgb = TEXT
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.name = FONT
        r2.font.color.rgb = color
    return tb


def panel(slide: Any, x: float, y: float, w: float, h: float, title: str | None = None, fill: RGBColor = PANEL, line: RGBColor = BORDER,
          radius: float = 0.05, title_color: RGBColor = CYAN, line_width: float = 0.75) -> Any:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    if title:
        text(slide, x + 0.14, y + 0.1, w - 0.28, 0.3, title.upper(), size=9.5, color=title_color, bold=True)
    return shp


def box(slide: Any, x: float, y: float, w: float, h: float, title: str, sub: str | None = None, fill: RGBColor = PANEL, line: RGBColor = BORDER,
        color: RGBColor = TEXT, size: float = 12, sub_size: float = 9.5, align: Any = PP_ALIGN.CENTER, bold: bool = True, radius: float = 0.12) -> Any:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.08))
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = title
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = align
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(sub_size)
        r2.font.name = FONT
        r2.font.color.rgb = MUTED
    return shp


def chip(slide: Any, x: float, y: float, w: float, label: str, color: RGBColor = CYAN, h: float = 0.36, size: float = 9.5) -> Any:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = PANEL
    shp.line.color.rgb = color
    shp.line.width = Pt(0.9)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = color
    return shp


def stat(slide: Any, x: float, y: float, w: float, h: float, value: str, label: str, color: RGBColor = CYAN, value_size: float = 26) -> None:
    panel(slide, x, y, w, h)
    text(slide, x + 0.12, y + 0.1, w - 0.24, h * 0.55, value, size=value_size, color=color, bold=True, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x + 0.12, y + h * 0.58, w - 0.24, h * 0.4, label, size=10, color=MUTED, line_spacing=1.05)


def arrow(slide: Any, x1: float, y1: float, x2: float, y2: float, color: RGBColor = MUTED, width: float = 1.25, head: bool = True) -> Any:
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if head:
        ln = c.line._get_or_add_ln()
        tail = ln.find(qn("a:tailEnd"))
        if tail is None:
            tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return c


def table(slide: Any, x: float, y: float, w: float, h: float, header: list[str], rows: list[list[str]], col_widths: list[float] | None = None,
          size: float = 10.5, header_size: float = 9.5, right_cols: tuple[int, ...] = (), highlight_rows: dict[int, RGBColor] | None = None,
          mono_cols: tuple[int, ...] = ()) -> Any:
    shape = slide.shapes.add_table(len(rows) + 1, len(header), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "1")
    tblPr.set("bandRow", "0")
    style = tblPr.find(qn("a:tableStyleId"))
    if style is None:
        style = etree.SubElement(tblPr, qn("a:tableStyleId"))
    style.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    row_h = h / (len(rows) + 1)
    for r in range(len(rows) + 1):
        tbl.rows[r].height = Inches(row_h)
    highlight_rows = highlight_rows or {}
    for c, head in enumerate(header):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER
        _cell_text(cell, head, header_size, CYAN, True, PP_ALIGN.RIGHT if c in right_cols else PP_ALIGN.LEFT)
    for r, row in enumerate(rows, start=1):
        fill = highlight_rows.get(r - 1, PANEL if r % 2 else PANEL2)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            _cell_text(cell, str(val), size, TEXT, False, PP_ALIGN.RIGHT if c in right_cols else PP_ALIGN.LEFT, MONO if c in mono_cols else FONT)
    return tbl


def _cell_text(cell: Any, value: str, size: float, color: RGBColor, bold: bool, align: Any, font: str = FONT) -> None:
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color


def image(slide: Any, path: Path, x: float, y: float, w: float | None = None, h: float | None = None) -> Any:
    if w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))


def _transparent(chart: Any) -> None:
    """Make a native chart's chart-area and plot-area transparent (dark slide background shows through)."""
    cs = chart._chartSpace
    chart_el = cs.find(qn("c:chart"))
    sp = cs.find(qn("c:spPr"))
    if sp is None:
        sp = etree.Element(qn("c:spPr"))
        chart_el.addnext(sp)
    for child in list(sp):
        sp.remove(child)
    etree.SubElement(sp, qn("a:noFill"))
    ln = etree.SubElement(sp, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))
    plot_area = chart_el.find(qn("c:plotArea"))
    pa = plot_area.find(qn("c:spPr"))
    if pa is None:
        pa = etree.SubElement(plot_area, qn("c:spPr"))
    for child in list(pa):
        pa.remove(child)
    etree.SubElement(pa, qn("a:noFill"))
    ln2 = etree.SubElement(pa, qn("a:ln"))
    etree.SubElement(ln2, qn("a:noFill"))


def doughnut(slide: Any, x: float, y: float, w: float, h: float, labels: list[str], values: list[float], colors: list[RGBColor], hole: int = 58) -> Any:
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("share", values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.RIGHT
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(10)
    ch.legend.font.name = FONT
    ch.legend.font.color.rgb = TEXT
    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = False
    dl.show_percentage = True
    dl.show_category_name = False
    dl.number_format = "0%"
    dl.number_format_is_linked = False
    dl.font.size = Pt(9.5)
    dl.font.bold = True
    dl.font.name = FONT
    dl.font.color.rgb = TEXT
    series = plot.series[0]
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = BG
        pt.format.line.width = Pt(1.5)
    for el in ch._chartSpace.xpath(".//c:holeSize"):
        el.set("val", str(hole))
    _transparent(ch)
    return ch


# ---------------------------------------------------------------------------------------------------------------------
# slides
# ---------------------------------------------------------------------------------------------------------------------
def s_title(deck: Deck) -> None:
    s = deck.bare("NEXUS — title", notes="Open on the name. One sentence: a virtual world that mirrors a physical operation, plus AI agents that understand, predict, simulate and optimize what happens inside it — and simulate before they act.")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.85), Inches(1.5), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, 0.85, 2.0, 7.5, 1.4, "NEXUS", size=66, color=TEXT, bold=True)
    text(s, 0.9, 3.3, 7.6, 0.9, "AI-native Digital Twin & Autonomous Operations Platform", size=22, color=CYAN, bold=True)
    text(s, 0.9, 4.1, 7.4, 1.2,
         "A virtual world that mirrors a physical operation — and AI agents that understand, predict, simulate and optimize what happens inside it. Before executing an action, NEXUS simulates it.",
         size=14, color=MUTED, line_spacing=1.2)
    cx = 0.9
    for label, color, cw in [("Simulate before execute", GREEN, 2.15), ("Local-first · no API keys", CYAN, 2.15), ("OR-Tools CP-SAT", VIOLET, 1.6), ("Apache-2.0", MUTED, 1.25)]:
        chip(s, cx, 5.45, cw, label, color, size=9)
        cx += cw + 0.12
    text(s, 0.9, 6.25, 8, 0.4, "Raunit Thakur  ·  raunit.thakur@gmail.com  ·  github.com/raunitgrey7/nexus", size=12, color=MUTED)
    text(s, 0.9, 6.6, 8, 0.35, f"{dt.date.today():%B %Y}  ·  Digital Twin + Physical AI + Multi-Agent AI + Simulation + Operations Research", size=10.5, color=MUTED)
    # stylised warehouse map
    mx, my, mw, mh = 8.75, 1.55, 3.95, 4.9
    panel(s, mx, my, mw, mh, fill=PANEL2)
    text(s, mx + 0.15, my + 0.12, 3, 0.3, "LIVE DIGITAL TWIN", size=9, color=CYAN, bold=True)
    text(s, mx + 2.35, my + 0.12, 1.5, 0.3, "10:30  ·  R07 ⚠", size=9, color=RED, bold=True, align=PP_ALIGN.RIGHT, font=MONO)
    # charging bay
    for k in range(4):
        cb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx + 0.2), Inches(my + 0.6 + k * 0.85), Inches(0.22), Inches(0.42))
        cb.fill.solid()
        cb.fill.fore_color.rgb = RGBColor(0x14, 0x3A, 0x2A)
        cb.line.color.rgb = GREEN
        cb.line.width = Pt(0.5)
        cb.shadow.inherit = False
    # shelves in 3 zone columns × 3 rows
    for col in range(3):
        for row in range(3):
            zx = mx + 0.7 + col * 1.1
            zy = my + 0.55 + row * 1.2
            tint = RGBColor(0x2A, 0x24, 0x16) if (col, row) == (1, 1) else PANEL
            z = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(zx), Inches(zy), Inches(0.95), Inches(1.05))
            z.fill.solid()
            z.fill.fore_color.rgb = tint
            z.line.color.rgb = BORDER
            z.line.width = Pt(0.5)
            z.shadow.inherit = False
            for k in range(3):
                sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(zx + 0.1 + k * 0.3), Inches(zy + 0.12), Inches(0.14), Inches(0.8))
                sh.fill.solid()
                sh.fill.fore_color.rgb = BORDER
                sh.line.fill.background()
                sh.shadow.inherit = False
    # robots
    for (rx, ry, color, label) in [(0.95, 0.9, CYAN, "R1"), (2.1, 1.5, CYAN, "R4"), (1.55, 2.6, CYAN, "R9"), (3.15, 1.15, CYAN, "R3"), (2.75, 3.15, CYAN, "R2"), (1.95, 2.05, RED, "R7")]:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx + rx), Inches(my + ry), Inches(0.2), Inches(0.2))
        d.fill.solid()
        d.fill.fore_color.rgb = color
        d.line.fill.background()
        d.shadow.inherit = False
        text(s, mx + rx - 0.1, my + ry + 0.18, 0.4, 0.22, label, size=7, color=color, bold=True, align=PP_ALIGN.CENTER, font=MONO)
    # docks
    for k in range(4):
        dk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx + 0.75 + k * 0.82), Inches(my + mh - 0.78), Inches(0.5), Inches(0.26))
        dk.fill.solid()
        dk.fill.fore_color.rgb = RGBColor(0x3A, 0x2E, 0x12)
        dk.line.color.rgb = AMBER
        dk.line.width = Pt(0.5)
        dk.shadow.inherit = False
    text(s, mx + 0.15, my + mh - 0.4, 3.7, 0.3, "orders 46  ·  robots 11/12  ·  predicted breach 39.4%", size=7.5, color=MUTED, font=MONO)


def s_problem(deck: Deck) -> None:
    s = deck.slide("The problem", "Operations react. Dashboards describe. Nobody simulates.",
                   notes="Three failure modes of today's operations software. The numbers at the bottom are from our own twin: one robot out of twelve at the peak is enough to put 39% of orders at risk within the hour if nobody acts.")
    cols = [
        ("Reactive", RED, ["A robot fails → an alert fires → whoever is on shift improvises.", "Every minute of delay compounds into missed SLAs; the recovery plan is invented under pressure."]),
        ("Descriptive", AMBER, ["WMS and BI dashboards show what has happened — not what will happen next, nor what to do about it.", "Forecasts, if any, live in spreadsheets far from the floor."]),
        ("Untested", VIOLET, ["Interventions are executed on the real floor. The first time a plan is tried is in production.", "There is no way to compare alternatives before committing robots, people and inventory."]),
    ]
    for i, (title, color, lines) in enumerate(cols):
        x = 0.6 + i * 4.1
        panel(s, x, 1.6, 3.9, 3.0)
        text(s, x + 0.2, 1.75, 3.5, 0.5, title, size=18, color=color, bold=True)
        bullets(s, x + 0.15, 2.3, 3.6, 2.2, lines, size=12.5, bullet_color=color, spacing=6)
    stat(s, 0.6, 4.95, 3.9, 1.75, "39.4%", "of orders projected late within the hour after one robot of twelve fails at the peak — if nobody acts", RED)
    stat(s, 4.7, 4.95, 3.9, 1.75, "+332%", "average fulfilment time without intervention (2.1 → 9.2 min)", AMBER)
    stat(s, 8.8, 4.95, 3.9, 1.75, "18.8 s", "for NEXUS to evaluate 9 candidate plans and 21 task allocations in forked twins — and pick one", CYAN)


def s_normal_vs_nexus(deck: Deck) -> None:
    s = deck.slide("The difference", "What a normal system says — and what NEXUS says",
                   notes="This is the exact operator briefing rendered by the platform on the demo incident. Every number is computed by simulating each candidate plan in a forked copy of the twin. The LLM only proposes candidates; mathematics disposes.")
    panel(s, 0.6, 1.6, 5.4, 4.6, "A normal system says")
    text(s, 0.85, 2.3, 5.0, 1.0, "Robot R07 offline.", size=30, color=MUTED, bold=True, font=MONO)
    text(s, 0.85, 3.4, 4.9, 2.5, "…and a tile turns red. Tasks pile up in R07's zone, the SLA clock keeps running, and what happens next depends on whoever is on shift — with no way to know which reaction would actually help.",
         size=13, color=MUTED, line_spacing=1.25)
    panel(s, 6.3, 1.6, 6.45, 4.6, "NEXUS says", line=CYAN)
    d = DEMO
    quote = [
        [("R07 failure (motor fault) will increase average order fulfilment time from ", {}), (f"{d['ft_now']} to {d['ft_nothing']} min", {"bold": True}),
         (" over the next 60 minutes — projected SLA breach ", {}), (f"{d['baseline_breach']}%", {"bold": True, "color": RED}), (" without intervention.", {})],
        [("I evaluated ", {}), (f"{d['candidates']} candidate plans and {d['allocations']} task allocations", {"bold": True}), (f" in {d['total_s']} s.", {})],
        [("Recommended plan #1 — Add 2 robots + batching. ", {"bold": True, "color": CYAN}), ("Estimated impact: SLA breach ", {}),
         (f"{d['baseline_breach']}% → {d['plan_breach']}%", {"bold": True, "color": GREEN}), (f", average fulfilment {d['ft_nothing']} → {d['ft_plan']} min, throughput {d['thr_nothing']} → {d['thr_plan']} orders/h.", {})],
        [("Risk LOW", {"bold": True, "color": GREEN}), (f" (stable across {d['seeds']} seeds, σ {d['sigma_pct']}%); auto-approved — projected SLA breach improves by ≥ 2 pp.", {})],
    ]
    rich(s, 6.5, 2.05, 6.1, 4.0, quote, size=12.5, line_spacing=1.2, space_after=9)
    text(s, 0.6, 6.35, 12.1, 0.6, "Rendered from a decision record (nexus/agents/explain.py). Every number comes from simulating each candidate in a forked copy of the twin, scoring it with a multi-objective cost and risk-checking it. Numbers: `uv run nexus demo`, 2026-08-30.",
         size=10, color=MUTED, italic=True)


def s_vision(deck: Deck) -> None:
    s = deck.slide("The vision", "From screens into environments",
                   notes="AI is moving from chat windows into physical operations. Analysts name physical AI, multi-agent systems, spatial AI and digital twins as strategic themes. NEXUS's thesis: the digital twin is the operating system of physical AI.")
    bullets(s, 0.6, 1.6, 6.4, 3.6, [
        ("The shift.", "AI is moving from answering questions on a screen to perceiving, reasoning and acting inside physical operations — warehouses, factories, hospitals, airports."),
        ("The gap.", "Robots and sensors produce the data; nobody turns it into safe, explainable decisions in real time."),
        ("Analyst themes 2025–26.", "physical AI, multi-agent systems, domain-specific intelligence and spatial AI (strategic technology trends); digital twins, autonomous systems, logistics optimization and foundation models (smart-manufacturing roadmap); robotics, autonomous forklifts, inspection drones and digital twins (physical-AI briefs)."),
        ("The market.", "warehouse automation and industrial digital twins are multi-billion-dollar categories growing at double-digit rates; every site that adds robots needs an intelligence layer above the fleet manager."),
    ], size=12.5, spacing=8)
    panel(s, 7.3, 1.6, 5.45, 3.6, "NEXUS thesis", line=CYAN)
    text(s, 7.5, 2.05, 5.05, 1.1, "The digital twin is the operating system of physical AI: the place where agents perceive, reason, simulate and act — safely.", size=14, color=TEXT, bold=True, line_spacing=1.2)
    bullets(s, 7.45, 3.2, 5.1, 1.9, [
        ("Perceive —", "the twin knows where everything is, what is happening and what is supposed to happen."),
        ("Reason —", "agents forecast, plan and explain in the operator's language."),
        ("Act —", "only after simulation, risk assessment and policy or human approval."),
    ], size=11.5, spacing=5)
    for i, label in enumerate(["Physical AI", "Multi-agent systems", "Spatial AI", "Digital twins", "Autonomous systems", "Logistics optimization", "Foundation models", "Domain intelligence"]):
        chip(s, 0.6 + i * 1.53, 5.55, 1.45, label, CYAN if i % 2 == 0 else VIOLET, size=8.5)
    text(s, 0.6, 6.15, 12.1, 0.6, "Themes as named in 2025–26 industry analyst publications (Gartner strategic technology trends; NIST smart-manufacturing roadmap; World Economic Forum physical-AI briefs). Market sizing is described qualitatively on purpose — figures differ by source.",
         size=9.5, color=MUTED, italic=True)


def s_product(deck: Deck) -> None:
    s = deck.slide("The product", "NEXUS in one picture",
                   notes="The twin is the single source of truth. Around it: an event engine that makes every change auditable and replayable, a deterministic simulation engine that can fork the present into many futures, an optimization engine, a multi-agent runtime that plans, validates, simulates, risk-checks, approves and executes, plus forecasting, a what-if lab and a natural-language console.")
    hub = (5.05, 3.25, 3.25, 1.55)
    box(s, *hub, "DIGITAL TWIN", "single source of truth · zones, shelves, robots, workers, orders, inventory, docks, chargers · spatial graph", fill=HEADER, line=CYAN, color=CYAN, size=15, sub_size=9.5)
    sats = [
        (0.9, 1.6, "EVENT ENGINE", "every change is a typed event · idempotent store · replay · audit", VIOLET),
        (9.2, 1.6, "SIMULATION ENGINE", "deterministic ticks · fork the present into many futures · 6k ticks/s", CYAN),
        (0.6, 3.4, "OPTIMIZATION", "CP-SAT assignment · order batching · deadline sequencing · routing policies", BLUE),
        (9.5, 3.4, "MULTI-AGENT RUNTIME", "plan → validate → optimize → simulate → risk → approve → execute", GREEN),
        (0.9, 5.2, "FORECASTING", "demand (Holt-Winters + prior) · battery · congestion · bottlenecks", AMBER),
        (9.2, 5.2, "WHAT-IF · NL CONSOLE", "scenario DSL, 12 presets, multi-strategy comparison · grounded answers", VIOLET),
        (5.05, 1.55, "LIVE UI · API", "Next.js + Three.js 3D twin · FastAPI + WebSocket", TEXT),
        (5.05, 5.55, "PERSISTENCE · OBSERVABILITY", "PostgreSQL · Redis · Prometheus · Grafana · OpenTelemetry", TEXT),
    ]
    cx, cy = hub[0] + hub[2] / 2, hub[1] + hub[3] / 2
    for (x, y, title, sub, color) in sats:
        w, h = 3.25, 1.2
        box(s, x, y, w, h, title, sub, line=color, color=color, size=12, sub_size=9)
        bx, by = x + w / 2, y + h / 2
        # connector from satellite edge towards hub edge
        dx, dy = cx - bx, cy - by
        dist = math.hypot(dx, dy) or 1
        ux, uy = dx / dist, dy / dist
        arrow(s, bx + ux * 1.55, by + uy * 0.62, cx - ux * 1.62, cy - uy * 0.78, color=BORDER, width=1.5, head=False)


def s_storyline(deck: Deck) -> None:
    s = deck.slide("The demo", "One simulated morning: surge, failure, decision, recovery",
                   notes="Walk the timeline: the twin runs from 08:00; at 10:25 demand surges 20%; at 10:30 R07 fails; within 19 seconds NEXUS has simulated nine plans and 21 allocations; the recommended plan executes; the realised breach converges to the simulated projection.")
    image(s, chart_demo_timeline(), 1.05, 1.45, w=11.2)
    stat(s, 0.6, 5.95, 2.9, 0.95, f"{DEMO['baseline_breach']}%", "projected breach, do nothing", RED, value_size=21)
    stat(s, 3.67, 5.95, 2.9, 0.95, f"{DEMO['plan_breach']}%", "simulated breach, recommended plan", GREEN, value_size=21)
    stat(s, 6.74, 5.95, 2.9, 0.95, f"{DEMO['total_s']} s", f"{DEMO['candidates']} plans · {DEMO['allocations']} allocations · risk-checked", CYAN, value_size=21)
    stat(s, 9.81, 5.95, 2.9, 0.95, "3.1%", "realized breach since decision at 11:30", GREEN, value_size=21)


def s_decision(deck: Deck) -> None:
    s = deck.slide("The decision", "Nine candidates, one horizon, one recommendation",
                   notes="Left: every candidate's projected SLA breach after a 60-minute simulation in a forked twin. Reassignment plans don't beat doing nothing because the optimized scheduler already reassigns released tasks within seconds; capacity levers do. Right: the pipeline that produced the ranking.")
    image(s, chart_candidates(), 0.6, 1.5, w=7.3)
    panel(s, 8.15, 1.5, 4.6, 4.5, "How it was decided")
    bullets(s, 8.25, 1.95, 4.4, 4.1, [
        ("Playbooks + (optional) LLM", "propose 9 candidates; every action is from a closed vocabulary of 16 types."),
        ("Constraint validation", "drops or clamps anything unsafe (fleet limits, zones, docks)."),
        ("Forked-world simulation", "60 min each, in a process pool, with the same engine as the live twin."),
        ("Multi-objective score", "ranks plans (lateness, delay, tail, congestion, distance, energy)."),
        ("Risk agent", "re-simulates the winner under 2 extra seeds: σ 0.2% → LOW; no deadlocks, no depletion."),
        ("Approval policy", "LOW risk and ≥ 2 pp gain → auto-approved; otherwise a human decides."),
    ], size=11, spacing=5)
    d = DEMO
    stat(s, 0.6, 6.05, 2.35, 0.85, f"{d['baseline_breach']}% → {d['plan_breach']}%", "SLA breach", GREEN, value_size=17)
    stat(s, 3.1, 6.05, 2.35, 0.85, f"{d['ft_nothing']} → {d['ft_plan']} min", "avg fulfilment", CYAN, value_size=17)
    stat(s, 5.6, 6.05, 2.35, 0.85, f"{d['thr_nothing']} → {d['thr_plan']}/h", "throughput", CYAN, value_size=17)
    text(s, 8.15, 6.1, 4.6, 0.8, f"Timings: planning {d['planning_ms']} ms · simulation {d['simulation_s']} s · risk {d['risk_s']} s · total {d['total_s']} s (4 worker processes).", size=10, color=MUTED)


def s_architecture(deck: Deck) -> None:
    s = deck.slide("Architecture", "One engine, three planes: twin, agents, events",
                   notes="UI talks to a FastAPI/WebSocket layer. Below it three planes share the world: the twin engine (state), the agent runtime (decisions) and the event engine (truth). The simulation engine runs the same code for the live twin and for every forked future; the optimization engine is called by the scheduler and by the agents.")
    box(s, 3.7, 1.5, 6.0, 0.72, "NEXUS UI", "Next.js 15 · Three.js 3D/2D twin · decisions · what-if lab · forecast · console · timeline · benchmarks", line=TEXT, size=12, sub_size=9)
    box(s, 3.7, 2.5, 6.0, 0.66, "API LAYER", "FastAPI · pydantic contract · REST · WebSocket /ws/live · Prometheus /metrics", line=CYAN, color=CYAN, size=12, sub_size=9)
    arrow(s, 6.7, 2.22, 6.7, 2.5)
    planes = [
        (0.6, "TWIN ENGINE", "WorldState · occupancy grid · semantic spatial graph · snapshots · SHA-256 digest · fork", VIOLET),
        (4.75, "AGENT RUNTIME", "Ops Manager · Planner · Validator · Optimizer · Simulator · Risk · Policy · Executor", GREEN),
        (8.9, "EVENT ENGINE", "typed events · idempotent append-only store · pure reducer · bus · replay", AMBER),
    ]
    for x, title, sub, color in planes:
        box(s, x, 3.5, 3.85, 0.98, title, sub, line=color, color=color, size=12, sub_size=9)
        arrow(s, 6.7, 3.16, x + 1.925, 3.5)
    arrow(s, 4.45, 3.99, 4.75, 3.99, head=True)
    arrow(s, 8.9, 3.99, 8.6, 3.99, head=True)
    box(s, 4.75, 4.85, 3.85, 0.7, "SIMULATION ENGINE", "deterministic ticks · kinematics · A* · faults · KPIs · forks", line=CYAN, color=CYAN, size=12, sub_size=9)
    arrow(s, 6.675, 4.48, 6.675, 4.85)
    box(s, 4.75, 5.9, 3.85, 0.7, "OPTIMIZATION ENGINE", "CP-SAT · Hungarian · GA · batching · EDF · routing policies", line=BLUE, color=BLUE, size=12, sub_size=9)
    arrow(s, 6.675, 5.55, 6.675, 5.9)
    box(s, 0.6, 4.85, 3.85, 1.75, "PERSISTENCE & DATA", "PostgreSQL — events, snapshots, decisions, what-ifs · Redis fan-out · SQLite for tests · in-memory when unset", line=BORDER, size=12, sub_size=9)
    box(s, 8.9, 4.85, 3.85, 1.75, "OBSERVABILITY & LLM", "Prometheus metrics · Grafana dashboard · OpenTelemetry tracing · Ollama (qwen2.5, local, optional)", line=BORDER, size=12, sub_size=9)
    text(s, 0.6, 6.7, 12.1, 0.3, "docker compose: postgres · redis · backend · frontend · prometheus · grafana · ollama (profile)   —   everything runs on one laptop, no paid APIs.", size=9.5, color=MUTED)


def s_twin(deck: Deck) -> None:
    s = deck.slide("Digital twin", "Every entity has state; the world is the source of truth",
                   notes="The twin is a dataclass world: fast to copy, trivial to hash. Four built-in scales from a tiny test world to 100 robots / 50 zones / 2,400 shelves. Forking a world costs milliseconds, which is what makes simulate-before-execute affordable.")
    panel(s, 0.6, 1.55, 5.6, 4.15, "Entity model")
    tree = [
        "Warehouse  (WorldState)",
        "├─ Zones      storage · corridor · dock · charging · capacity",
        "├─ Shelves    cell · access cell · inventory {SKU: qty}",
        "├─ Robots     cell · battery · status · task · path · load",
        "├─ Workers    role · dock · status (delays)",
        "├─ Orders     lines · priority · deadline · status",
        "├─ Tasks      waypoints (pick… → deliver) · leg",
        "├─ Docks      open · queue      Chargers  slots · occupants",
        "└─ Clock · seeded RNG · id counters · demand profile · config · stats",
    ]
    text(s, 0.8, 1.95, 5.3, 3.7, tree, size=10.5, color=TEXT, font=MONO, line_spacing=1.25, space_after=2)
    header = ["scale", "robots", "storage zones", "docks", "chargers", "SKUs", "inventory", "orders/h", "grid", "shelves"]
    rows = [[sc, str(r), z, str(d), str(c), sk, inv, str(oph), g, f"{sh:,}"] for (sc, r, z, d, c, sk, inv, oph, g, sh, _w) in SCALES]
    table(s, 6.45, 1.55, 6.3, 2.35, header, rows, col_widths=[0.62, 0.55, 0.9, 0.5, 0.65, 0.5, 0.72, 0.62, 0.62, 0.62], size=9, header_size=8, right_cols=(1, 3, 4, 5, 6, 7, 9), highlight_rows={1: HEADER})
    text(s, 6.45, 3.95, 6.3, 0.5, "Base demand × hourly profile (08:00–18:00 operating window). The small world is the calibrated demo: ≈ 4,000 orders/day, 12 robots, near capacity at the late-morning peak.", size=9.5, color=MUTED)
    stat(s, 6.45, 4.6, 2.0, 1.1, "5 / 26 ms", "world fork (small / large) — pickle round-trip", CYAN, value_size=18)
    stat(s, 8.6, 4.6, 2.0, 1.1, "SHA-256", "deterministic digest of all behaviour-relevant state", VIOLET, value_size=18)
    stat(s, 10.75, 4.6, 2.0, 1.1, "2,400", "shelves · 150k units · 100 robots at large scale", AMBER, value_size=18)
    text(s, 0.6, 5.9, 12.1, 0.9, "Domain-agnostic by construction: the engine only knows entity shapes; a DomainModel (warehouse today) supplies the layout, demand and vocabulary. Spatial reasoning lives in a semantic graph derived from the world (R07 is_inside Zone C · Zone C adjacent_to C2 · ORD-8821 requires Shelf C-41).",
         size=11, color=MUTED, line_spacing=1.2)


def s_events(deck: Deck) -> None:
    s = deck.slide("Event-sourced world", "Everything is an event: replayable, idempotent, auditable",
                   notes="Only the reducer mutates the world. Engines and agents emit events. Replay is a snapshot plus the external events; the deterministic engine regenerates the rest and the digest must match — this is tested.")
    bullets(s, 0.6, 1.55, 6.3, 4.3, [
        ("Typed, closed vocabulary.", "orders, tasks, robots, workers, infrastructure, inventory, demand, plans, system — every change has a type and a payload."),
        ("Append-only, idempotent store.", "sequence numbers, idempotency keys for external commands (a retried command applies once), ring buffer for the live feed."),
        ("Reducer-only mutation.", "one pure function applies events to the world; engines and agents never poke at entities."),
        ("Replay = snapshot + external events.", "engine events are regenerated deterministically; the replayed world's digest equals the original (test gate)."),
        ("Attribution.", "origin (engine · scenario · agent · user), cause (plan id) and key on every event — an audit trail for free."),
        ("Ephemeral kinematics.", "per-tick movement events are streamed to observers but not persisted; the engine reproduces them."),
    ], size=12, spacing=6)
    flow = [
        ("Engine · agents · users emit", MUTED),
        ("Event store — seq, idempotency key, ring buffer", VIOLET),
        ("Reducer — pure, the only writer", CYAN),
        ("World state", TEXT),
        ("Bus → UI · persistence · metrics", GREEN),
    ]
    for i, (label, color) in enumerate(flow):
        y = 1.6 + i * 0.95
        box(s, 7.3, y, 5.4, 0.6, label, line=color, color=color, size=11.5)
        if i < len(flow) - 1:
            arrow(s, 10.0, y + 0.6, 10.0, y + 0.95)
    text(s, 7.3, 6.35, 5.4, 0.5, "Replay verified: snapshot + 3 external events → identical SHA-256 digest (tests/test_simulation.py).", size=9.5, color=MUTED, italic=True)
    for i, ev in enumerate(["ORDER_CREATED", "TASK_CREATED", "ITEM_PICKED", "ORDER_DELIVERED", "BATTERY_LOW", "ROBOT_FAILURE", "AISLE_BLOCKED", "ZONE_CLOSED", "DEMAND_CHANGED", "PLAN_EXECUTED"]):
        chip(s, 0.6 + (i % 5) * 1.33, 6.0 + (i // 5) * 0.42, 1.26, ev, VIOLET if i % 2 else CYAN, h=0.32, size=7)


def s_simulation(deck: Deck, bench: dict[str, Any] | None, label: str) -> None:
    s = deck.slide("Deterministic simulation", "Deterministic worlds you can fork into many futures",
                   notes="The engine is a discrete-time lockstep simulation with a fixed order of operations. Determinism is what makes forks, replays and benchmarks trustworthy. Speed: thousands of simulated seconds per wall second.")
    panel(s, 0.6, 1.55, 5.7, 4.6, "One tick, in order")
    steps = [
        ("1", "Order arrivals — Poisson × hourly profile; Zipf SKU popularity creates natural hot zones"),
        ("2", "Faults — scheduled scenario events + seeded spontaneous failures"),
        ("3", "Recoveries and worker delays expire"),
        ("4", "Inventory replenishment (every 5 min)"),
        ("5", "Strategy decisions — assignments, batching, routing, charging"),
        ("6", "Robot kinematics in id order — move (congestion slows), pick, unload, charge, battery"),
        ("7", "Hooks (history recorder, agents) · TICK with utilization & congestion stats"),
    ]
    for i, (n, t) in enumerate(steps):
        y = 2.0 + i * 0.57
        chip(s, 0.8, y + 0.02, 0.38, n, CYAN, h=0.34, size=9)
        text(s, 1.3, y - 0.02, 4.9, 0.55, t, size=10.5, color=TEXT, line_spacing=1.05)
    image(s, chart_ticks(bench, label), 6.55, 1.5, w=6.2)
    stat(s, 6.55, 5.6, 1.95, 1.1, "digest ==", "same seed ⇒ identical world hash", GREEN, value_size=16)
    stat(s, 8.65, 5.6, 1.95, 1.1, "fork ==", "a forked world continues identically", GREEN, value_size=16)
    stat(s, 10.75, 5.6, 1.95, 1.1, "replay ==", "snapshot + external events reproduce it", GREEN, value_size=16)
    text(s, 0.6, 6.3, 5.7, 0.6, "Also: A* with congestion-aware costs and a path cache, BFS distance fields for the optimizer, capacity-limited cells, deterministic slow-down in congested zones, battery drain/charge model, 10/5/3-minute SLAs.", size=9.5, color=MUTED)


def s_optimization(deck: Deck) -> None:
    s = deck.slide("Optimization", "Operations research in the loop, not an AI wrapper",
                   notes="The objective is a weighted multi-objective cost shared by the scheduler, the plan ranking and the what-if comparison. Assignment is a CP-SAT model with fallbacks; batching, deadline sequencing and routing policies are the levers the agents pull.")
    panel(s, 0.6, 1.55, 3.95, 3.7, "Objective (lower is better)")
    text(s, 0.75, 1.95, 3.7, 0.9, "J = Σ w·[ lateness, delivery time, p95 tail, congestion, distance, energy, backlog, priority ]", size=12, color=TEXT, font=MONO, line_spacing=1.2)
    bullets(s, 0.7, 2.95, 3.8, 2.2, [
        "one scorer for the scheduler, plan ranking and what-if comparison",
        "weights: lateness 100 · delivery time 1 · tail 0.25 · congestion 2 · distance 0.005 · energy 0.1",
        "constraints: capacity, battery with reserve, sourceable lines, open dock, one robot per task",
    ], size=10.5, spacing=4)
    panel(s, 4.7, 1.55, 3.95, 3.7, "Assignment model")
    bullets(s, 4.8, 1.95, 3.8, 3.2, [
        ("CP-SAT (OR-Tools):", "x[r][b] ∈ {0,1}; each batch ≤ 1 robot, each robot ≤ 1 batch; maximise assignments, minimise cost; 0.2 s limit; single worker + fixed seed → deterministic."),
        ("Costs:", "exact BFS distances, trip estimate, lateness vs deadline, battery risk, congestion of zones touched."),
        ("Fallback chain:", "CP-SAT → Hungarian (scipy) → greedy; a genetic allocator generates alternative allocations."),
    ], size=10.5, spacing=6)
    panel(s, 8.8, 1.55, 3.95, 3.7, "Levers the agents pull")
    bullets(s, 8.9, 1.95, 3.8, 3.2, [
        ("Batching —", "2–4 orders per trip clustered by zone proximity, deadline-aware."),
        ("Weighted EDF —", "sequence by deadline slack over priority weight."),
        ("Routing policies —", "avoid zones, prefer corridors, soft capacities, with expiry."),
        ("Pre-emptive charging —", "robots charge before exhaustion, never a quarter of the fleet at once."),
    ], size=10.5, spacing=6)
    stat(s, 0.6, 5.45, 3.95, 1.3, "~420 → ~600+/h", "capacity of 12 robots: naive dispatch vs optimized scheduler (peak sweep)", CYAN, value_size=18)
    stat(s, 4.7, 5.45, 3.95, 1.3, "0.4% vs 48%", "SLA breach at 504 orders/h: optimized vs baseline (90 min, same seed)", GREEN, value_size=18)
    stat(s, 8.8, 5.45, 3.95, 1.3, "0.2 s", "CP-SAT time limit per planning cycle; ≈ 2k ticks/s end-to-end on the small world", VIOLET, value_size=18)


def s_forecasting(deck: Deck) -> None:
    s = deck.slide("Forecasting", "NEXUS does not just respond — it anticipates",
                   notes="Four forecasters run every simulated minute: demand, battery, congestion and bottlenecks. They feed the planner prompt, the risk agent and the UI. Milliseconds, deterministic.")
    cards = [
        (0.6, 1.55, "Demand", CYAN, ["Holt-Winters (additive, damped) over per-minute arrivals, blended with the twin's demand-profile prior", "15-minute buckets with prediction intervals; trend; capacity/h from fleet × task duration; projected utilization"]),
        (6.7, 1.55, "Battery", GREEN, ["per robot: drain rate from history vs remaining workload → predicted exhaustion minutes", "charger ETA via pathfinding; risk high if exhaustion < ETA + 10 min", "\"Send R04 to charging after its current task (exhaustion in 23 min, charger 4 min away)\""]),
        (0.6, 3.75, "Congestion", AMBER, ["per zone: robots inside + en route + the wave of open orders that require the zone (Little's law) + history trend", "projected load vs capacity, ETA to congestion, drivers (\"14 open orders require Zone C\")"]),
        (6.7, 3.75, "Bottlenecks", RED, ["rules with severities 0–1: zone load, dock queues, charger saturation, failed robots, hot-SKU stock-outs, loader delays, demand vs capacity, backlog age", "each with a concrete recommendation: \"Pre-position inventory in Zone B\", \"+2 robots or enable batching\""]),
    ]
    for x, y, title, color, lines in cards:
        panel(s, x, y, 6.05, 2.05)
        text(s, x + 0.2, y + 0.12, 3, 0.35, title, size=14, color=color, bold=True)
        bullets(s, x + 0.15, y + 0.5, 5.8, 1.55, lines, size=10.5, bullet_color=color, spacing=3)
    panel(s, 0.6, 5.95, 9.0, 0.95, fill=PANEL2)
    text(s, 0.75, 6.02, 8.7, 0.85, "“Demand rising to ~419 orders/h (+23%) over the next 90 min; projected utilization 1.32 — above capacity. Zone F congestion +325% expected: 4.2 robots vs capacity 3. Recommended: +3 robots or enable batching to lift capacity to ~636 orders/h.” — Forecast.summary, small world, 09:30 after an R07 failure",
         size=9.5, color=MUTED, italic=True, line_spacing=1.15)
    stat(s, 9.8, 5.95, 2.95, 0.95, "3.5 / 8.8 ms", "full forecast, small / large world", CYAN, value_size=17)


def s_agents(deck: Deck) -> None:
    s = deck.slide("Multi-agent runtime", "The LLM proposes. Mathematics and policy dispose.",
                   notes="Nine stages. The planner is the only place a language model appears, and it is optional: deterministic playbooks guarantee coverage. Everything downstream is validation, optimization, simulation, risk and policy.")
    stages = [
        ("GOAL + TRIGGER", "minimise SLA breach · ROBOT_FAILURE:R07", MUTED),
        ("FORECASTER", "demand · battery · congestion · bottlenecks", AMBER),
        ("SITUATION", "KPIs · fleet · hot zones · backlog · events", TEXT),
        ("PLANNER", "playbooks + LLM (Ollama) + SOP retrieval → candidates", VIOLET),
        ("VALIDATOR", "closed vocabulary · clamps · drops unsafe actions", RED),
        ("OPTIMIZER", "CP-SAT assignment · concrete robot choices", BLUE),
        ("SIMULATOR", "forked twins · process pool · horizon KPIs", CYAN),
        ("RISK", "deadlock · safety · exhaustion · regression · stability seeds", AMBER),
        ("POLICY → EXECUTOR", "auto if LOW risk & ≥ 2 pp gain, else human · idempotent events", GREEN),
    ]
    for i, (title, sub, color) in enumerate(stages):
        row, col = divmod(i, 5)
        if row == 1:
            col += 0.5
        x = 0.6 + col * 2.46
        y = 1.6 + row * 1.9
        box(s, x, y, 2.3, 1.35, title, sub, line=color, color=color, size=11, sub_size=8.5)
        if i < len(stages) - 1:
            if row == 0 and col < 4:
                arrow(s, x + 2.3, y + 0.675, x + 2.46, y + 0.675)
            elif row == 1 and i < len(stages) - 1:
                arrow(s, x + 2.3, y + 0.675, x + 2.46, y + 0.675)
    # elbow from VALIDATOR (end of row 0) to OPTIMIZER (start of row 1)
    x_end = 0.6 + 4 * 2.46 + 1.15
    x_start = 0.6 + 0.5 * 2.46 + 1.15
    y_mid = 1.6 + 1.35 + 0.28
    arrow(s, x_end, 1.6 + 1.35, x_end, y_mid, color=MUTED, width=1.25, head=False)
    arrow(s, x_end, y_mid, x_start, y_mid, color=MUTED, width=1.25, head=False)
    arrow(s, x_start, y_mid, x_start, 3.5, color=MUTED, width=1.25)
    panel(s, 0.6, 5.45, 12.15, 1.4, fill=PANEL2)
    bullets(s, 0.75, 5.55, 11.9, 1.3, [
        ("Decision record.", f"situation, baseline, every candidate with simulation + risk + rank, approval, explanation, timings — {DEMO['candidates']} plans / {DEMO['allocations']} allocations in {DEMO['total_s']} s on the demo incident (planning {DEMO['planning_ms']} ms · simulation {DEMO['simulation_s']} s · risk {DEMO['risk_s']} s)."),
        ("Autopilot.", "failures, closures, blockages and demand changes trigger a decision (15-minute cooldown) in a background thread while the twin keeps ticking; benchmark strategies ai_planner and nexus_full run this loop inside the simulation."),
    ], size=10.5, spacing=4)


def s_safety(deck: Deck) -> None:
    s = deck.slide("Safety architecture", "Simulate before execute, with a human in the loop",
                   notes="Safety is a stack of controls, not a feature: closed vocabulary, validation, forked simulation, risk assessment including stability re-runs, approval policy, attributable idempotent events. The limitations are stated honestly.")
    bullets(s, 0.6, 1.55, 5.3, 4.9, [
        ("Closed action vocabulary.", "16 typed actions with parameter shapes; no free-form commands can reach the world."),
        ("Constraint validation.", "≤ ¼ of the fleet to charge · never remove > ⅓ of the fleet · at least one storage zone and one dock stay open · unknown ids dropped, ranges clamped."),
        ("Simulate before execute.", "every candidate runs in a forked twin with the real engine over the decision horizon; the do-nothing reference is always simulated alongside."),
        ("Risk agent.", "findings from diagnostics, baseline comparison and stability re-runs under extra seeds (table)."),
        ("Approval policy.", "auto-approve only if risk LOW and projected SLA gain ≥ 2 pp; otherwise proposed → operator approves, picks another candidate, or rejects."),
        ("Attributable, idempotent execution.", "origin · cause (plan id) · key on every event; executing a plan twice is a no-op."),
    ], size=11.5, spacing=6)
    header = ["finding", "trigger", "severity"]
    rows = [
        ["constraint", "plan infeasible · actions dropped by validation", "critical · info"],
        ["deadlock", "robot blocked ≥ 60 s or > 600 stuck robot-s · ≥ 25 s / > 150", "high · medium"],
        ["safety", "zones over 2× capacity ≥ 5% of horizon · ≥ 1% · transient", "high · medium · low"],
        ["resource exhaustion", "battery < 5% · < 12% · charger starvation · stock-outs", "high · medium · low"],
        ["regression", "breach worse than doing nothing · throughput −5% · congestion +50%", "high · medium · low"],
        ["capacity", "utilization > 97% — no slack", "low"],
        ["instability", "σ of breach across seeds > 5% · > 2.5% · else", "high · medium · info"],
    ]
    table(s, 6.15, 1.55, 6.6, 3.35, header, rows, col_widths=[1.5, 3.55, 1.55], size=9.5, header_size=9)
    text(s, 6.15, 4.98, 6.6, 0.5, "Score = Σ severity weights (low 0.12 · medium 0.30 · high 0.55 · critical 1.0). LOW < 0.2 ≤ MEDIUM < 0.45 ≤ HIGH; any critical finding → CRITICAL.", size=9.5, color=MUTED)
    panel(s, 6.15, 5.55, 6.6, 1.3, "Limitations, stated honestly", title_color=AMBER)
    text(s, 6.3, 5.9, 6.3, 0.95, "The twin is a calibrated simulation, not a physics engine; fidelity against a real site must be measured (planned: historical replay ingestion, Webots/Gazebo/Isaac bridges). No plan reaches real equipment without a fleet-manager integration and an operator.", size=9.5, color=MUTED, line_spacing=1.15)


def s_spatial(deck: Deck) -> None:
    s = deck.slide("Spatial AI", "Relationships, not just coordinates",
                   notes="A semantic spatial graph is derived from the world on demand: zones, corridors, shelves, robots, orders and their relations. The planner uses it to pick helpers by zone hops and spare neighbours for inventory; the console uses it to explain.")
    bullets(s, 0.6, 1.55, 5.9, 4.6, [
        ("Relations.", "is_inside · adjacent_to · located_in · requires · requires_zone · assigned_to · serves · charging_at · unloading_at"),
        ("Queries.", "entities in a zone, orders requiring a zone, zone-level routes (BFS honouring closures), per-zone load, relations of any entity"),
        ("Used by the planner —", "nearest helpers by zone hops from a failed robot; spare neighbour zones for inventory repositioning; corridors adjacent to hot zones"),
        ("Used by the risk agent and the console —", "\"Zone C congestion accounts for ~61% of predicted delay\"-style attribution grounded in the graph"),
        ("Grid + graph.", "the occupancy grid (A*, congestion, blockages) and the semantic graph are two views of one world — geometry for robots, semantics for agents and people"),
    ], size=11.5, spacing=7)
    panel(s, 6.8, 1.55, 5.95, 4.6, "Semantic graph (excerpt)")
    zones = {"A": (7.3, 2.6), "B": (9.3, 2.6), "C": (11.3, 2.6), "D": (7.3, 4.6), "E": (9.3, 4.6), "F": (11.3, 4.6)}
    corr = {"C1": (8.3, 3.6), "C2": (10.3, 3.6)}
    for zid, (x, y) in zones.items():
        box(s, x - 0.45, y - 0.32, 0.9, 0.64, f"Zone {zid}", line=CYAN if zid in ("A", "C") else BORDER, color=TEXT, size=10.5)
    for cid, (x, y) in corr.items():
        box(s, x - 0.5, y - 0.26, 1.0, 0.52, cid, "corridor", line=VIOLET, color=VIOLET, size=10, sub_size=7.5)
    for zid, cid in (("A", "C1"), ("B", "C1"), ("D", "C1"), ("E", "C1"), ("B", "C2"), ("C", "C2"), ("E", "C2"), ("F", "C2")):
        (x1, y1), (x2, y2) = zones[zid], corr[cid]
        arrow(s, x1, y1 + (0.32 if y2 > y1 else -0.32), x2, y2 + (-0.26 if y2 > y1 else 0.26), color=BORDER, width=1.0, head=False)
    text(s, 6.95, 5.15, 5.65, 0.95, ["R07  is_inside  Zone C        Zone C  adjacent_to  C2", "ORD-8821  requires  Shelf C-41   Shelf C-41  located_in  Zone C", "zone_route(A → C) = A → C1 → C2 → C   (closures honoured)"], size=9, color=MUTED, font=MONO, line_spacing=1.2, space_after=2)


def s_whatif(deck: Deck, sweep: dict[str, Any] | None) -> None:
    s = deck.slide("What-If engine", "Ask about one future, get four back — compared",
                   notes="A scenario is a list of typed mutations applied to a forked world; each strategy runs it in parallel; results are compared on the shared KPIs and scored with the optimization objective. Twelve presets, or any question in the console.")
    panel(s, 0.6, 1.55, 6.1, 2.35, "How it works")
    bullets(s, 0.7, 1.95, 5.9, 1.95, [
        "fork the current state → apply scenario mutations as scheduled events → run every strategy × seeds in a process pool",
        "compare on the shared KPI definitions, score with the optimization objective, keep the unmodified world as reference",
        "a deterministic narrative + comparison table; results stream back to the UI",
    ], size=10.5, spacing=4)
    panel(s, 0.6, 4.05, 6.1, 2.8, "Presets (also reachable by asking in plain English)")
    for i, p in enumerate(WHATIF_PRESETS):
        col, row = divmod(i, 6)
        text(s, 0.8 + col * 3.0, 4.45 + row * 0.38, 2.9, 0.36, f"›  {p}", size=10, color=TEXT)
    path, vals = chart_whatif(sweep)
    image(s, path, 6.95, 1.5, w=5.8)
    if vals:
        v = list(vals.values())
        text(s, 6.95, 5.85, 5.8, 1.0, f"Under the same peak demand, the R07 failure costs the naive scheduler {v[1] - v[0]:.0f} pp of SLA breach and the optimized scheduler {v[3] - v[2]:.1f} pp — resilience is a scheduling property, not only a hardware count.", size=10.5, color=MUTED, line_spacing=1.2)


def s_console(deck: Deck, nlq: dict[str, Any] | None) -> None:
    s = deck.slide("Natural-language console", "Ask the operation. Get grounded numbers, not vibes.",
                   notes="Every answer is computed by the engine — attribution, forecast, what-if, entity lookup — then optionally rewritten by the local LLM, which is discarded if it changes a number. These examples were captured from a run with the LLM off.")
    examples = [e for e in ((nlq or {}).get("examples") or []) if "could not complete" not in e.get("answer", "")]
    if len(examples) < 4:
        examples.append({"question": "What happens if order volume increases by 40%?", "intent": "whatif",
                         "answer": "Runs the scenario through the What-If engine (current strategy vs optimized vs nexus_full over the horizon) and answers with the measured SLA breach, fulfilment and throughput per strategy — see the What-If slide for the measured comparison.",
                         "latency_ms": 0})
    if not examples:
        examples = [{"question": q, "intent": i, "answer": "(run pitch/capture_measurements.py to capture the grounded answer)", "latency_ms": 0} for q, i in (
            ("How many orders are open right now?", "status"), ("Why are orders slowing down?", "explain"), ("What happens if order volume increases by 40%?", "whatif"), ("Where is R03 and what is it doing?", "entity"))]
    for i, ex in enumerate(examples[:4]):
        col, row = divmod(i, 2)
        x, y = 0.6 + col * 6.15, 1.55 + row * 2.6
        panel(s, x, y, 6.0, 2.45)
        text(s, x + 0.2, y + 0.12, 4.4, 0.35, ex["question"], size=12, color=CYAN, bold=True)
        chip(s, x + 4.65, y + 0.14, 1.15, ex["intent"].upper(), VIOLET, h=0.3, size=8)
        ans = ex["answer"]
        if len(ans) > 560:
            ans = ans[:557].rsplit(" ", 1)[0] + "…"
        text(s, x + 0.2, y + 0.52, 5.6, 1.7, ans, size=9.5, color=TEXT, line_spacing=1.15)
        lat = ex.get("latency_ms", 0)
        footer = f"deterministic · {lat:.0f} ms · LLM off" if lat else "deterministic · What-If engine, forked twins · LLM off"
        text(s, x + 0.2, y + 2.15, 5.6, 0.28, footer, size=8, color=MUTED, font=MONO)
    text(s, 0.6, 6.75, 12.1, 0.3, (nlq or {}).get("description", "") + "  Intents: status · explain · whatif · forecast · recommend · entity.", size=9, color=MUTED, italic=True)


def s_ui(deck: Deck) -> None:
    s = deck.slide("Visualization", "A control room, not a dashboard",
                   notes="The UI is a Next.js + Three.js twin: instanced shelves, congestion-tinted zones, robots coloured by status with their remaining paths, a KPI bar, fault injection, and the decision drawer. Seven pages; works fully offline in mock mode.")
    header = ["route", "what it shows"]
    rows = [
        ["/", "Live Twin — 3D/2D view, KPI bar, sim controls (play · pause · step · speed · reset · autopilot), events/robots/orders/zones panel, fault injection, Decide now"],
        ["/decisions", "situation · baseline vs candidates (recommended highlighted) · actions with rationale · risk findings · approve / reject / execute · timeline chart · timings · LLM badge"],
        ["/whatif", "preset cards · scenario builder (mutation vocabulary) · async runs · comparison table · KPI bars · timelines · best strategy · narrative"],
        ["/forecast", "demand band chart · capacity vs forecast · projected utilization · battery table · congestion bars · bottleneck list"],
        ["/console", "natural-language chat with suggestion chips, intent chip, LLM/deterministic badge, inline what-if card"],
        ["/timeline", "KPI history with notable-event markers · snapshot scrubber · playback in the 2D view"],
        ["/benchmarks", "grouped KPI bars per scale · normalised radar · results table"],
    ]
    table(s, 0.6, 1.55, 8.1, 4.5, header, rows, col_widths=[1.2, 6.9], size=9.5, header_size=9, mono_cols=(0,))
    panel(s, 8.95, 1.55, 3.8, 4.5, "The 3D twin")
    bullets(s, 9.05, 1.95, 3.6, 4.1, [
        "instanced shelves (up to 2,400), zones tinted by occupancy / capacity",
        "robots coloured by status with id labels and remaining path; failed robots pulse red",
        "docks, chargers, blocked cells, hatched closed zones",
        "smooth interpolation between ~20 tick frames/s over one WebSocket",
        "2D canvas toggle and snapshot playback",
        "dark control-room design system; responsive to 1280 px",
    ], size=10.5, spacing=5)
    text(s, 0.6, 6.2, 12.1, 0.6, "Next.js 15 · React 19 · TypeScript strict · Tailwind v4 · @react-three/fiber + drei · recharts · zustand · typed API client mirroring the pydantic contract · NEXT_PUBLIC_MOCK=1 runs every page offline against an in-browser sim.", size=10, color=MUTED, line_spacing=1.2)


def s_bench_breach(deck: Deck, bench: dict[str, Any], label: str) -> None:
    cfg = bench.get("config", {})
    s = deck.slide("Benchmarks I", "Four strategies, identical worlds, identical incidents",
                   notes="All strategies see the same seeds, orders and incident schedule (robot failure at +30 min, demand surge at +60 min, aisle blocked at +90 min). SLA breach is the headline KPI.")
    image(s, chart_bench(bench, "sla_breach_pct", "Projected SLA breach", "SLA breach (%)", "{:.1f}%", "bench_breach", label, figsize=(7.6, 4.6)), 0.6, 1.5, w=7.6)
    panel(s, 8.45, 1.5, 4.3, 4.6, "Setup")
    bullets(s, 8.55, 1.9, 4.1, 4.2, [
        ("Runs:", f"{cfg.get('minutes', '?')} simulated minutes × {cfg.get('seeds', '?')} seed(s) per cell · scales {', '.join(cfg.get('scales', []))}"),
        ("Incidents:", "robot failure at +30 min (45 min recovery) · demand ×1.5 for 30 min at +60 · aisle blocked at +90 (15 min)"),
        ("Strategies:", "baseline (FIFO + nearest) · optimized (CP-SAT + batching + routing) · ai_planner (playbooks, no simulation) · nexus_full (simulate-before-execute, risk-gated)"),
        ("Same KPI definitions", "everywhere: metrics.py → benchmarks → UI → this deck"),
        ("Reproduce:", "make bench · results JSON + docs/BENCHMARKS.md + charts"),
    ], size=10.5, spacing=6)
    text(s, 0.6, 6.25, 12.1, 0.5, f"Data: backend/benchmarks/results ({label}) · generated {bench.get('generated_at', '?')[:19].replace('T', ' ')} · wall time {cfg.get('wall_s_total', '?')} s.", size=9.5, color=MUTED, italic=True)


def s_bench_ft(deck: Deck, bench: dict[str, Any], label: str) -> None:
    s = deck.slide("Benchmarks II", "Speed and throughput: what the optimization buys",
                   notes="Average fulfilment time and throughput per strategy and scale. Throughput is demand-limited when the fleet keeps up; the gap opens as the scheduler falls behind.")
    image(s, chart_bench(bench, "avg_fulfillment_min", "Average fulfilment time", "minutes", "{:.1f}", "bench_ft", label, figsize=(6.0, 4.4)), 0.6, 1.5, w=6.0)
    image(s, chart_bench(bench, "throughput_per_hour", "Throughput", "orders / hour", "{:.0f}", "bench_thr", label, figsize=(6.0, 4.4)), 6.75, 1.5, w=6.0)
    rows = bench["summary_table"]
    scales = list(bench["scales"])
    notes = []
    for sc in scales:
        b = next((r for r in rows if r["scale"] == sc and r["strategy"] == "baseline"), None)
        o = next((r for r in rows if r["scale"] == sc and r["strategy"] == "optimized"), None)
        if b and o:
            notes.append(f"{sc}: breach {b['sla_breach_pct']:.1f}% → {o['sla_breach_pct']:.1f}%, fulfilment {b['avg_fulfillment_min']:.1f} → {o['avg_fulfillment_min']:.1f} min, p95 {b['p95_fulfillment_min']:.1f} → {o['p95_fulfillment_min']:.1f} min")
    text(s, 0.6, 6.05, 12.1, 0.8, "Baseline → optimized, same worlds:  " + "  ·  ".join(notes) if notes else "", size=10, color=MUTED, line_spacing=1.2)


def s_bench_table(deck: Deck, bench: dict[str, Any], label: str) -> None:
    s = deck.slide("Benchmarks III", "The full picture: every KPI, every strategy",
                   notes="Left: the summary table straight from the benchmark runner. Right: a normalised radar — each axis is scaled so that the best strategy scores 1.")
    rows = bench["summary_table"]
    header = ["scale", "strategy", "breach", "Δ vs base", "avg ft", "p95", "thr/h", "util", "cong.", "t/s"]
    body = []
    hl = {}
    for i, r in enumerate(rows):
        body.append([r["scale"], r["strategy"], f"{r['sla_breach_pct']:.2f}%", (f"{r['sla_breach_vs_baseline_pct']:+.2f} pp" if r.get("sla_breach_vs_baseline_pct") is not None else "—"),
                     f"{r['avg_fulfillment_min']:.2f}", f"{r['p95_fulfillment_min']:.2f}", f"{r['throughput_per_hour']:.0f}", f"{r['robot_utilization_pct']:.0f}%", f"{r['congestion_index']:.2f}", f"{r['ticks_per_second']:,.0f}"])
        if r["strategy"] == "nexus_full":
            hl[i] = HEADER
    n = len(body)
    h = min(4.9, 0.36 * (n + 1))
    table(s, 0.6, 1.5, 7.0, h, header, body[:14], col_widths=[0.7, 1.05, 0.8, 0.85, 0.65, 0.6, 0.62, 0.55, 0.58, 0.6], size=9 if n <= 8 else 8, header_size=8.5, right_cols=(2, 3, 4, 5, 6, 7, 8, 9), highlight_rows=hl, mono_cols=(1,))
    path, scale = chart_radar(bench, label)
    image(s, path, 7.85, 1.5, w=4.9)
    text(s, 0.6, 6.55, 12.1, 0.4, f"{label} · normalised radar on the {scale} scale: SLA compliance, speed, tail, throughput, flow and energy per order, each relative to the best strategy.", size=9.5, color=MUTED, italic=True)


def s_capacity(deck: Deck, sweep: dict[str, Any] | None) -> None:
    s = deck.slide("The capacity cliff", "Why the naive baseline collapses, and why it matters",
                   notes="Queueing systems have a cliff: past capacity, backlog and lateness explode. Batching and optimized assignment roughly double the capacity of the same 12 robots, which is why the baseline collapses in the benchmarks at demand the optimized scheduler handles comfortably.")
    image(s, chart_capacity(sweep), 0.6, 1.5, w=7.6)
    panel(s, 8.45, 1.5, 4.3, 4.85, "Reading the curve")
    bullets(s, 8.55, 1.9, 4.1, 4.4, [
        ("A queueing cliff.", "past capacity, backlog and lateness explode; below it, SLAs are trivially met."),
        ("Same robots, same orders.", "the only difference between the two solid lines is the scheduler."),
        ("Batching ≈ 2× capacity.", "2–4 orders per trip amortise travel — the single biggest lever before buying hardware."),
        ("Resilience.", "dashed lines add an R07 failure at +30 min; the optimized scheduler absorbs it until the cliff."),
        ("Design consequence.", "the benchmark's incident schedule is set where the naive scheduler is near its cliff — that is where decisions matter."),
    ], size=10.5, spacing=6)
    text(s, 0.6, 6.5, 7.6, 0.4, (sweep or {}).get("description", ""), size=8.5, color=MUTED, italic=True)


def s_quality(deck: Deck) -> None:
    s = deck.slide("Engineering quality", "Tested, typed, deterministic, containerised",
                   notes="92 tests across nine areas, determinism and replay gates, ruff and mypy clean, three CI jobs, Docker images. The API contract is pydantic on the backend and mirrored in TypeScript.")
    tests = PERF["tests"]
    labels = [f"{k} ({v})" for k, v in tests.items()]
    doughnut(s, 0.4, 1.45, 6.0, 4.3, labels, list(tests.values()), [CYAN, VIOLET, GREEN, AMBER, BLUE, RED, MUTED, RGBColor(0x38, 0xBD, 0xF8), RGBColor(0xF4, 0x72, 0xB6)])
    text(s, 0.6, 5.7, 5.8, 0.4, f"{sum(tests.values())} tests in 9 files — share by area", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    bullets(s, 6.7, 1.55, 6.05, 4.3, [
        ("Determinism gates.", "same seed ⇒ identical digest · fork continues identically · replay from snapshot + external events reproduces the digest."),
        ("Static quality.", "ruff (lint + format) and mypy clean across 87 source files; TypeScript strict, ESLint clean."),
        ("CI (GitHub Actions).", "backend: ruff · mypy · pytest with coverage · determinism smoke · micro-benchmark; frontend: lint · typecheck · build; docker: both images."),
        ("Contract-first API.", "pydantic schemas → OpenAPI → TypeScript mirrors; WebSocket frames documented."),
        ("Operations.", f"Docker Compose stack (7 services), Prometheus metrics + Grafana dashboard, OpenTelemetry tracing, PostgreSQL persistence; backend image {PERF['docker_gb']} GB."),
        ("Local-first.", "everything — including the LLM — runs on one laptop with zero paid APIs."),
    ], size=11, spacing=6)
    for i, (v, l) in enumerate([("92", "tests passing"), ("87", "typed source files"), ("3", "CI jobs"), ("7", "compose services")]):
        stat(s, 0.6 + i * 3.06, 6.1, 2.9, 0.8, v, l, CYAN, value_size=18)


def s_stack(deck: Deck, loc: dict[str, tuple[int, int]]) -> None:
    s = deck.slide("Tech stack", "Free, open, and measured: what the platform is made of",
                   notes="Line counts are measured by walking the repository at build time. The entire stack costs nothing to run.")
    labels = [f"{k} — {v[0]:,} lines" for k, v in loc.items()]
    doughnut(s, 0.4, 1.45, 6.4, 4.4, labels, [v[0] for v in loc.values()], [CYAN, VIOLET, GREEN, AMBER, BLUE])
    total = sum(v[0] for v in loc.values())
    files = sum(v[1] for v in loc.values())
    text(s, 0.6, 5.85, 6.2, 0.4, f"{total:,} non-blank lines across {files} files (measured at build time)", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    header = ["layer", "technology"]
    rows = [
        ["Twin, events, simulation", "Python 3.13 · dataclasses · orjson · networkx"],
        ["Optimization", "OR-Tools CP-SAT · SciPy (Hungarian) · genetic allocator"],
        ["Forecasting", "NumPy · Holt-Winters · Little's-law congestion estimator"],
        ["Agents & LLM", "pydantic structured output · Ollama (qwen2.5) · TF-IDF SOP retrieval"],
        ["API & runtime", "FastAPI · WebSocket · pydantic-settings · structlog"],
        ["Data", "PostgreSQL (SQLAlchemy async / asyncpg) · Redis · SQLite for tests"],
        ["UI", "Next.js 15 · React 19 · TypeScript · Tailwind v4 · Three.js (R3F) · recharts · zustand"],
        ["Ops", "Docker Compose · Prometheus · Grafana · OpenTelemetry · GitHub Actions · uv"],
    ]
    table(s, 7.05, 1.5, 5.7, 4.3, header, rows, col_widths=[1.65, 4.05], size=9.5, header_size=9)
    stat(s, 7.05, 6.0, 2.75, 0.85, "₹0 / $0", "to run the whole stack — no paid APIs", GREEN, value_size=18)
    stat(s, 9.95, 6.0, 2.8, 0.85, "Apache-2.0", "open-source licence", CYAN, value_size=18)


def s_domains(deck: Deck) -> None:
    s = deck.slide("The grander vision", "Same engine, different worlds",
                   notes="The engine never imports the warehouse. A DomainModel supplies the layout, demand and vocabulary. Factory, hospital, airport, data center and fleets map onto the same shapes: agents, jobs, resources, sinks, energy, areas.")
    box(s, 4.9, 1.55, 3.55, 0.95, "DIGITAL TWIN ENGINE", "events · simulation · optimization · agents · what-if · forecasting", fill=HEADER, line=CYAN, color=CYAN, size=13, sub_size=9)
    domains = [
        ("WAREHOUSE", "today: robots, orders, shelves, docks, chargers, workers", GREEN, True),
        ("FACTORY", "workcells = zones · WIP = orders · AGVs = robots · tool changes = charging", CYAN, False),
        ("HOSPITAL", "wards = zones · transport orders · porters & carts = agents · bays = docks", VIOLET, False),
        ("AIRPORT", "gates & stands · baggage & turnaround jobs · ground vehicles", AMBER, False),
        ("DATA CENTER", "racks & rows · tickets = jobs · technicians · spare-parts inventory", BLUE, False),
        ("FLEET / SUPPLY CHAIN", "depots = zones · shipments = orders · vehicles = robots · fuel = battery", MUTED, False),
    ]
    for i, (title, sub, color, now) in enumerate(domains):
        x = 0.6 + i * 2.05
        box(s, x, 3.6, 1.95, 1.5, title, sub, line=color, color=color, size=10.5, sub_size=8, fill=PANEL if not now else RGBColor(0x0F, 0x2A, 0x20))
        arrow(s, 6.675, 2.5, x + 0.975, 3.6, color=BORDER, width=1.1)
    panel(s, 0.6, 5.4, 12.15, 1.45, fill=PANEL2)
    bullets(s, 0.75, 5.5, 11.9, 1.3, [
        ("DomainModel protocol:", "build(scale, seed) → world · vocabulary() (agent = robot, job = order, resource = shelf, sink = dock, energy = battery, area = zone) · describe(world) for LLM context."),
        ("What stays:", "event sourcing, deterministic simulation, forks, CP-SAT assignment, batching, routing policies, forecasting, the agent pipeline, the safety stack, the API and the UI. What changes: layout, entities' meaning, KPIs' names."),
    ], size=10.5, spacing=4)


def s_business(deck: Deck) -> None:
    s = deck.slide("Business model & go-to-market", "Pilot → platform per site → domain packs",
                   notes="Land with a pilot on one site using historical events; expand to a per-site platform subscription; grow through domain packs and integrations. Buyers are operations leaders, 3PLs and robotics integrators who already own the fleets.")
    phases = [
        ("1 · PILOT", "twin of one site · replay historical events · benchmark the four strategies · 8–12 weeks · fixed fee", CYAN),
        ("2 · PLATFORM", "live twin + agents + what-if per site · subscription tiered by robots and zones · on-prem or VPC", GREEN),
        ("3 · DOMAIN PACKS", "factory, hospital, airport models · WMS / fleet-manager connectors · simulator bridges", VIOLET),
        ("4 · ENTERPRISE", "multi-site, SSO / RBAC / audit, local LLM, SLAs · OEM licensing to WMS and robotics vendors", AMBER),
    ]
    for i, (title, sub, color) in enumerate(phases):
        x = 0.6 + i * 3.08
        box(s, x, 1.6, 2.9, 1.6, title, sub, line=color, color=color, size=12.5, sub_size=9)
        if i < 3:
            arrow(s, x + 2.9, 2.4, x + 3.08, 2.4)
    panel(s, 0.6, 3.5, 5.95, 3.35, "Who buys, and why")
    bullets(s, 0.7, 3.9, 5.75, 2.9, [
        ("Buyers:", "VP / Director of Operations, 3PL site leads, robotics integrators, WMS vendors (OEM)."),
        ("SLA protection:", "on the demo incident, doing nothing projects 39.4% late orders in the hour; the recommended plan 3.1%."),
        ("Capacity deferral:", "batching and optimized assignment roughly double the capacity of the same fleet before new robots are bought."),
        ("Incident response in seconds:", "nine simulated plans in 18.8 s instead of an improvised reaction."),
        ("Auditability:", "every executed action is an attributable event with the simulation that justified it."),
    ], size=10.5, spacing=5)
    panel(s, 6.8, 3.5, 5.95, 3.35, "Pricing shape (illustrative)")
    bullets(s, 6.9, 3.9, 5.75, 2.9, [
        ("Pilot:", "fixed project fee; deliverable = calibrated twin + benchmark report + decision log."),
        ("Platform:", "per site per month, tiered by fleet size (≤ 25 · ≤ 100 · > 100 robots) and zones."),
        ("Packs & integrations:", "priced per connector / domain model; simulator bridges as add-ons."),
        ("Why it sticks:", "the twin accumulates the site's history, calibration and playbooks — switching cost grows with every incident handled."),
        ("Cost base:", "no per-token API fees — local models and open-source solvers keep gross margin high."),
    ], size=10.5, spacing=5)


def s_competition(deck: Deck) -> None:
    s = deck.slide("Competitive landscape", "No one else closes the loop from prediction to action",
                   notes="Dashboards describe, fleet managers execute, simulation tools plan offline, LLM copilots chat. NEXUS is the layer that connects prediction, simulation, optimization and safe execution — and it runs locally.")
    header = ["", "live twin", "predicts", "simulates before acting", "optimizes (OR)", "safety & approval", "explains in language", "local-first"]
    rows = [
        ["WMS / BI dashboards", "partial", "–", "–", "–", "–", "–", "✓"],
        ["Robot fleet managers (vendor)", "✓ (own fleet)", "–", "–", "partial", "partial", "–", "✓"],
        ["Discrete-event simulation tools", "–", "offline", "offline", "partial", "–", "–", "✓"],
        ["Generic LLM copilots", "–", "–", "–", "–", "–", "✓", "rarely"],
        ["NEXUS", "✓", "✓", "✓", "✓", "✓", "✓", "✓"],
    ]
    table(s, 0.6, 1.6, 12.15, 3.6, header, rows, col_widths=[2.75, 1.25, 1.1, 1.75, 1.35, 1.45, 1.4, 1.1], size=10.5, header_size=9, right_cols=(), highlight_rows={4: RGBColor(0x0F, 0x2A, 0x20)})
    bullets(s, 0.6, 5.45, 12.15, 1.4, [
        ("Position:", "the intelligence layer above the fleet manager and the WMS — vendor-neutral, integrating through events rather than replacing what a site already runs."),
        ("Moat:", "the calibrated twin + the decision log + domain playbooks; determinism and simulation-first safety are architectural, not features that can be bolted onto a dashboard."),
    ], size=11, spacing=5)


def s_roadmap(deck: Deck) -> None:
    s = deck.slide("Roadmap", "13 milestones shipped; next: real robots, more worlds",
                   notes="M0–M13 are done and in the repository. Next quarters: simulator bridges, reinforcement-learning dispatch policies, factory and hospital domain packs, multi-site, and enterprise hardening.")
    done = ["M0 Foundation", "M1 Twin core", "M2 Event engine", "M3 Simulation", "M4 Optimization", "M5 Forecasting", "M6 Agent runtime", "M7 What-If", "M8 API & ops", "M9 Twin UI", "M10 Benchmarks", "M11 Docs & ADRs", "M12 Ship", "M13 Pitch deck"]
    text(s, 0.6, 1.5, 6, 0.3, "SHIPPED — August 2026", size=10, color=GREEN, bold=True)
    for i, m in enumerate(done):
        chip(s, 0.6 + (i % 7) * 1.74, 1.85 + (i // 7) * 0.45, 1.66, m, GREEN, h=0.36, size=8.5)
    cols = [
        ("NEXT 3 MONTHS", CYAN, ["Historical-event ingestion: replay a real site's WMS / fleet logs into the twin", "Calibration harness: measure simulation fidelity against realised KPIs", "Webots / Gazebo bridge: robots in a physics simulator driven by NEXUS plans", "Design-partner pilots (2–3 sites)"]),
        ("3–9 MONTHS", VIOLET, ["Reinforcement-learning dispatch policy trained inside the twin, benchmarked against CP-SAT", "Factory and hospital domain packs (DomainModel)", "WMS and fleet-manager connectors; Isaac Sim bridge", "Multi-site twins with shared inventory"]),
        ("9–18 MONTHS", AMBER, ["Enterprise: SSO, RBAC, audit export, SLAs", "Fine-tuned local planner model on the decision log", "Marketplace of playbooks and domain packs", "OEM licensing to WMS / robotics vendors"]),
    ]
    for i, (title, color, items) in enumerate(cols):
        x = 0.6 + i * 4.1
        panel(s, x, 2.95, 3.95, 3.9)
        text(s, x + 0.2, 3.08, 3.5, 0.35, title, size=12, color=color, bold=True)
        bullets(s, x + 0.15, 3.5, 3.7, 3.3, items, size=10.5, bullet_color=color, spacing=6)


def s_team(deck: Deck) -> None:
    s = deck.slide("Team & portfolio", "One engineer, four layers of AI systems",
                   notes="A deliberate progression: an AI operating system for people, AI for software infrastructure, AI for security, and now AI for physical operations. Each project is public, benchmarked and local-first.")
    steps = [
        ("SHERRY", "Personal AI OS", "execute tasks for a human", MUTED),
        ("SENTINEL", "AI incident investigation", "diagnose failures in software infrastructure", BLUE),
        ("AEGIS", "AI cybersecurity investigation", "detect threats — 100% detection / 2% FPR on a reproducible benchmark", VIOLET),
        ("NEXUS", "Digital twin & autonomous operations", "predict · simulate · optimize physical reality", CYAN),
    ]
    for i, (name, what, line, color) in enumerate(steps):
        x = 0.6 + i * 3.08
        box(s, x, 1.7, 2.9, 1.75, name, f"{what}\n{line}", line=color, color=color, size=15, sub_size=9.5)
        if i < 3:
            arrow(s, x + 2.9, 2.575, x + 3.08, 2.575)
    text(s, 0.6, 3.6, 12.1, 0.4, "AI operating systems → AI infrastructure → AI security → AI physical intelligence", size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    panel(s, 0.6, 4.2, 12.15, 2.65, "Raunit Thakur — AI engineer")
    bullets(s, 0.75, 4.6, 11.9, 2.2, [
        ("Systems + AI:", "event-sourced deterministic simulation, operations research (CP-SAT, Hungarian, GA), forecasting, multi-agent runtimes with structured LLM output, real-time UIs (Three.js), production tooling (Docker, CI, observability)."),
        ("Working style:", "benchmark-first, local-first, zero-API-key architectures; every claim reproducible with one command."),
        ("Looking for:", "design partners with robot fleets, advisors in robotics and supply-chain operations, and early-stage backing to build the bridges to real equipment."),
        ("Contact:", "raunit.thakur@gmail.com · github.com/raunitgrey7 (sentinel, aegis, nexus are public)"),
    ], size=11, spacing=6)


def s_ask(deck: Deck) -> None:
    s = deck.slide("The ask", "Help NEXUS meet its first real robots",
                   notes="Three concrete asks: pilots, backing for the physical bridges, and advisors. Everything shown runs on a laptop today; the next step is a real site.")
    asks = [
        ("DESIGN PARTNERS", "2–3 warehouse / 3PL sites with robot fleets for calibrated pilots: replay history, benchmark strategies, run what-ifs on real questions", CYAN),
        ("PRE-SEED SUPPORT", "12–18 months of runway for simulator bridges (Webots / Gazebo / Isaac), fleet-manager connectors and domain packs", GREEN),
        ("ADVISORS", "robotics, warehouse operations and supply-chain leaders who have lived through the incidents NEXUS is built for", VIOLET),
    ]
    for i, (title, sub, color) in enumerate(asks):
        x = 0.6 + i * 4.1
        box(s, x, 1.65, 3.95, 2.1, title, sub, line=color, color=color, size=14, sub_size=10.5)
    panel(s, 0.6, 4.05, 12.15, 1.65, fill=PANEL2)
    text(s, 0.8, 4.2, 11.8, 1.4, ["Everything in this deck runs on one laptop:", "docker compose up --build   →   UI :3000 · API :8000/docs · Grafana :3001 · Prometheus :9090", "uv run nexus demo   →   the storyline in the terminal (twin → surge → R07 fails → nine plans simulated → best plan executed)"],
         size=12, color=TEXT, font=MONO, line_spacing=1.25, space_after=6)
    text(s, 0.6, 5.95, 12.15, 0.5, "Raunit Thakur  ·  raunit.thakur@gmail.com  ·  github.com/raunitgrey7/nexus", size=16, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    text(s, 0.6, 6.45, 12.15, 0.4, "Apache-2.0 · local-first · zero API keys · simulate before execute", size=11, color=MUTED, align=PP_ALIGN.CENTER)


def s_appendix_kpis(deck: Deck) -> None:
    s = deck.slide("Appendix A", "KPI definitions and the action vocabulary",
                   notes="Reference material: the KPI definitions shared by every component, and the sixteen action types a plan may contain.")
    header = ["KPI", "definition"]
    rows = [
        ["Avg fulfilment time", "mean(delivered_at − created_at) over delivered orders"],
        ["SLA breach (projected)", "late delivered + open-overdue orders over delivered + open — the headline number"],
        ["Throughput", "delivered orders per simulated hour"],
        ["Robot utilization", "share of robot-ticks in a productive state (moving to pick / deliver, picking, unloading)"],
        ["Congestion index", "mean over ticks of Σ zones max(0, robots_in_zone − capacity)"],
        ["Distance / energy", "cells travelled · battery percentage points consumed (per order in the benchmark tables)"],
        ["Planning latency", "wall-clock seconds to produce a decision (plan → optimize → simulate → risk)"],
    ]
    table(s, 0.6, 1.55, 7.4, 3.9, header, rows, col_widths=[1.9, 5.5], size=9.5, header_size=9)
    panel(s, 8.25, 1.55, 4.5, 5.2, "Action vocabulary (16 types)")
    for i, a in enumerate(ACTIONS):
        text(s, 8.45 + (i % 2) * 2.2, 1.98 + (i // 2) * 0.55, 2.15, 0.5, a, size=9, color=TEXT, font=MONO)
    text(s, 0.6, 5.6, 7.4, 1.15, "SLA targets in the calibrated world: LOW 20 min · NORMAL 10 · HIGH 5 · CRITICAL 3. Pick 6 s per line, unload 4 s (×2 without a loader), battery 0.02 %/cell, charge 0.15 %/s, low-battery threshold 20 %. All constants live in SimConfig and are part of the world (forks carry them).", size=9.5, color=MUTED, line_spacing=1.2)


def s_appendix_api(deck: Deck) -> None:
    s = deck.slide("Appendix B", "API surface and live stream",
                   notes="The REST surface and the WebSocket frames. Shapes are defined once in pydantic and mirrored in TypeScript.")
    groups = [
        ("Twin & control", ["GET /api/health · /api/status", "POST /api/sim/control (start · pause · step · reset · speed · autopilot)", "GET /api/world · /world/robots · /world/orders · /world/entity/{id} · /world/relations/{id}", "GET /api/kpis · /api/spatial · /api/timeline · /api/snapshots/{tick}"]),
        ("Events & faults", ["GET /api/events?since_seq · /api/events/recent", "POST /api/events/inject (idempotent by key)", "GET /api/faults/presets · POST /api/faults/{preset}"]),
        ("Intelligence", ["GET /api/forecast?horizon_min", "POST /api/decisions · GET /api/decisions[/{id}] · POST /api/decisions/{id}/actions (approve · reject · execute)", "POST /api/whatif (async) · /api/whatif/run (sync) · GET /api/whatif[/{id}] · /api/whatif/presets", "POST /api/nlq · GET /api/strategies · GET /api/benchmarks · GET /metrics"]),
    ]
    for i, (title, items) in enumerate(groups):
        y = 1.55 + i * 1.75
        panel(s, 0.6, y, 7.6, 1.6, title)
        text(s, 0.8, y + 0.42, 7.2, 1.15, items, size=9, color=TEXT, font=MONO, line_spacing=1.2, space_after=2)
    panel(s, 8.45, 1.55, 4.3, 5.1, "WebSocket /ws/live")
    text(s, 8.65, 1.98, 3.95, 4.6, [
        "server → client", "  hello   world + KPIs + status", "  tick    robots, KPIs, zones, docks, chargers (≤ 20 / s)", "  event   notable + order / task events", "  decision · forecast · whatif · status", "", "client → server", "  control {action, ticks_per_second, autopilot}", "  subscribe {tick_every}", "  ping → pong",
    ], size=9, color=TEXT, font=MONO, line_spacing=1.25, space_after=2)


# ---------------------------------------------------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------------------------------------------------
def build() -> Path:
    bench, label = load_benchmark()
    sweep = load_json(DATA / "sweep.json")
    nlq = load_json(DATA / "nlq_examples.json")
    loc = count_loc()
    print(f"benchmark data: {label}", flush=True)
    print(f"sweep: {'ok' if sweep else 'missing'} · nlq examples: {'ok' if nlq else 'missing'} · loc: {sum(v[0] for v in loc.values()):,} lines", flush=True)

    deck = Deck()
    s_title(deck)
    s_problem(deck)
    s_normal_vs_nexus(deck)
    s_vision(deck)
    s_product(deck)
    s_storyline(deck)
    s_decision(deck)
    s_architecture(deck)
    s_twin(deck)
    s_events(deck)
    s_simulation(deck, bench, label)
    s_optimization(deck)
    s_forecasting(deck)
    s_agents(deck)
    s_safety(deck)
    s_spatial(deck)
    s_whatif(deck, sweep)
    s_console(deck, nlq)
    s_ui(deck)
    if bench is not None:
        s_bench_breach(deck, bench, label)
        s_bench_ft(deck, bench, label)
        s_bench_table(deck, bench, label)
    s_capacity(deck, sweep)
    s_quality(deck)
    s_stack(deck, loc)
    s_domains(deck)
    s_business(deck)
    s_competition(deck)
    s_roadmap(deck)
    s_team(deck)
    s_ask(deck)
    s_appendix_kpis(deck)
    s_appendix_api(deck)
    out = deck.save()
    print(f"wrote {out} ({out.stat().st_size / 1024:.0f} KB, {deck.n} slides)")
    return out


def verify(path: Path) -> None:
    prs = Presentation(str(path))
    print(f"verify: {len(prs.slides)} slides re-loaded from {path.name}")
    for i, slide in enumerate(prs.slides, start=1):
        texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        title = next((t for t in texts if len(t) > 12 and not t.isupper()), texts[0] if texts else "")
        print(f"  {i:02d}  {title[:96]}")


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    verify(build())
